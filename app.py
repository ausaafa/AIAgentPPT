from flask import Flask, request, jsonify, send_from_directory
from werkzeug.utils import secure_filename
import subprocess
import tempfile
import shutil
from pathlib import Path
import os
import re
from datetime import datetime
from copy import deepcopy

import uuid  # ← ADD this import at the top of app.py if not present
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime

import cv2

from PIL import Image, ImageDraw, ImageFilter

import numpy as np
from PIL import Image, ImageDraw
from dotenv import load_dotenv
from PIL import Image, ImageDraw, ImageFont
import pytesseract
if os.name == "nt":
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
import cv2
import numpy as np
from openai import OpenAI
from docx import Document
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGBColor
from PyPDF2 import PdfReader

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
    from pptx.dml.color import RGBColor as PptRGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ---------------- CONFIG ----------------
app = Flask(__name__, static_folder="static", template_folder="templates")

UPLOAD_FOLDER = "uploads"
GENERATED_FOLDER = "generated"
ALLOWED_EXTENSIONS = {"txt", "docx", "pdf"}
ALLOWED_IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "bmp", "tiff", "gif"}
IMAGES_FOLDER = "images"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["GENERATED_FOLDER"] = GENERATED_FOLDER
app.config["IMAGES_FOLDER"] = IMAGES_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GENERATED_FOLDER, exist_ok=True)
os.makedirs(IMAGES_FOLDER, exist_ok=True)
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


"""
Drop-in replacement for the image translation section of app.py

Key fixes:
  1. Uses Tesseract --psm 11 (sparse text) for better multi-region detection
  2. Groups words into lines using the (block, par, line) key — then merges
     adjacent lines that belong to the same visual paragraph block
  3. Infers font size from the OCR bounding box height instead of hardcoding 14px
  4. Erases a generous padding around the source text before drawing overlay
  5. Finds the best-fit font size that actually fits within the box width
  6. Handles DPI: Pillow always works in pixel space, so no unit conversion needed
     for the PIL path; the docx blob replacement is also pixel-correct
"""

import os
import re
from datetime import datetime

import cv2
import numpy as np
from PIL import Image, ImageDraw, ImageFont

# Global job progress tracker
job_progress = {}
# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------
def _load_font(size: int) -> ImageFont.FreeTypeFont:
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf",   # Arial Bold (BEST on Windows)
        "C:/Windows/Fonts/calibrib.ttf",  # Calibri Bold
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    ]

    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except:
                continue

    return ImageFont.load_default()

def _fit_font_to_box(text: str, box_w: int, box_h: int,
                     initial_size: int) -> tuple[ImageFont.FreeTypeFont, int]:
    """
    Return a (font, size) whose text fits within box_w × box_h.
    Starts at initial_size and shrinks by 1pt until it fits (minimum 8pt).
    """
    size = max(initial_size, 8)
    for attempt_size in range(size, 7, -1):
        font = _load_font(attempt_size)
        # Use a throw-away draw to measure
        dummy_img = Image.new("RGB", (1, 1))
        dummy_draw = ImageDraw.Draw(dummy_img)
        bbox = dummy_draw.textbbox((0, 0), text, font=font)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
        if tw <= box_w and th <= box_h:
            return font, attempt_size
    return _load_font(8), 8


# ---------------------------------------------------------------------------
# Public API  (replaces the originals in app.py)
# ---------------------------------------------------------------------------

def _preprocess_for_ocr(pil_image: Image.Image) -> Image.Image:
    """
    Sharpen and lightly denoise so Tesseract reads bold/display text better.
    Returns a new PIL image; does NOT modify the original.
    """
    # Convert to OpenCV
    cv_img = cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)
 
    # Mild denoise (fast, preserves edges on bold text)
    cv_img = cv2.fastNlMeansDenoisingColored(cv_img, None, 5, 5, 7, 21)
 
    # Unsharp mask — brings out character edges
    blurred = cv2.GaussianBlur(cv_img, (0, 0), 2.0)
    cv_img = cv2.addWeighted(cv_img, 1.6, blurred, -0.6, 0)
 
    return Image.fromarray(cv2.cvtColor(cv_img, cv2.COLOR_BGR2RGB))
 
 
def _is_garbage_region(text: str) -> bool:
    """
    Return True if the OCR'd text looks like a misread / artifact.
    Heuristics:
      - Fewer than 2 real word characters
      - More than 40 % punctuation / symbols
      - All single characters separated by spaces (e.g. "H e S a e d")
      - Repeating character patterns
    """
    stripped = text.strip()
    if len(stripped) < 2:
        return True
 
    words = stripped.split()
 
    # All single-char "words" → likely broken character segmentation
    if len(words) > 1 and all(len(w) == 1 for w in words):
        return True
 
    alpha_count = sum(c.isalpha() for c in stripped)
    punct_count = sum(not c.isalnum() and not c.isspace() for c in stripped)
 
    # Fewer than 2 alpha characters total
    if alpha_count < 2:
        return True
 
    # More than 40 % punctuation
    if len(stripped) > 0 and punct_count / len(stripped) > 0.40:
        return True
 
    return False
 
 
def detect_text_regions(image_path: str):
    import pytesseract
 
    try:
        pil_image = Image.open(image_path).convert("RGB")
        image_cv = cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)
    except Exception as e:
        print(f"PIL could not read image {image_path}: {e}")
        return [], None
 
    # ← CHANGED: sharpen before OCR
    sharpened = _preprocess_for_ocr(pil_image)
 
    all_line_maps = {}
 
    try:
        ocr_data = pytesseract.image_to_data(
            sharpened,
            output_type=pytesseract.Output.DICT,
            lang="eng",
            config="--psm 6",
        )
    except Exception as e:
        print(f"OCR failed: {e}")
        return [], pil_image
 
    for i, word in enumerate(ocr_data["text"]):
        word = word.strip()
        if not word:
            continue
        try:
            conf = float(ocr_data["conf"][i])
        except Exception:
            conf = -1
        if conf < 50:                           # ← CHANGED: was 30
            continue
 
        key = (ocr_data["block_num"][i], ocr_data["par_num"][i], ocr_data["line_num"][i])
        x = ocr_data["left"][i]
        y = ocr_data["top"][i]
        w = ocr_data["width"][i]
        h = ocr_data["height"][i]
 
        if key not in all_line_maps:
            all_line_maps[key] = {"words": [], "xs": [], "ys": [],
                                  "rights": [], "bottoms": [], "heights": []}
        e = all_line_maps[key]
        e["words"].append(word)
        e["xs"].append(x);     e["ys"].append(y)
        e["rights"].append(x + w); e["bottoms"].append(y + h)
        e["heights"].append(h)
 
    raw_regions = []
    for entry in all_line_maps.values():
        text = " ".join(entry["words"]).strip()
        if _is_garbage_region(text):            # ← CHANGED: garbage filter
            print(f"  [OCR filter] skipped garbage region: {repr(text)}")
            continue
 
        min_x, min_y = min(entry["xs"]), min(entry["ys"])
        max_x, max_y = max(entry["rights"]), max(entry["bottoms"])
        median_h = float(np.median(entry["heights"])) if entry["heights"] else 14
 
        raw_regions.append({
            "text": text,
            "bbox": (min_x, min_y, max_x - min_x, max_y - min_y),
            "font_size_est": int(median_h * 0.85),
        })
 
    if not raw_regions:
        return [], pil_image
 
    raw_regions.sort(key=lambda r: (r["bbox"][1], r["bbox"][0]))
    deduped = []
    for reg in raw_regions:
        rx, ry, rw, rh = reg["bbox"]
        duplicate = False
        for kept in deduped:
            kx, ky, kw, kh = kept["bbox"]
            ix  = max(rx, kx); iy  = max(ry, ky)
            ix2 = min(rx+rw, kx+kw); iy2 = min(ry+rh, ky+kh)
            if ix2 > ix and iy2 > iy and rw*rh > 0:
                if (ix2-ix)*(iy2-iy) / (rw*rh) > 0.5:
                    duplicate = True; break
        if not duplicate:
            deduped.append(reg)
 
    return _merge_lines_into_blocks(deduped), pil_image

def _merge_lines_into_blocks(regions: list, max_v_gap: int = 8, max_x_diff: int = 30) -> list:
    """
    Merge single-line regions that are close vertically and share the same
    left-edge x into multi-line blocks.
    """
    if not regions:
        return regions

    regions = sorted(regions, key=lambda r: (r["bbox"][1], r["bbox"][0]))
    groups = [[regions[0]]]

    for reg in regions[1:]:
        rx, ry, rw, rh = reg["bbox"]
        last = groups[-1][-1]
        lx, ly, lw, lh = last["bbox"]
        vertical_gap = ry - (ly + lh)
        x_diff = abs(rx - lx)

        if vertical_gap <= max_v_gap and x_diff <= max_x_diff:
            groups[-1].append(reg)
        else:
            groups.append([reg])

    merged = []
    for group in groups:
        if len(group) == 1:
            merged.append(group[0])
            continue

        all_text = " ".join(g["text"] for g in group)
        min_x = min(g["bbox"][0] for g in group)
        min_y = min(g["bbox"][1] for g in group)
        max_x = max(g["bbox"][0] + g["bbox"][2] for g in group)
        max_y = max(g["bbox"][1] + g["bbox"][3] for g in group)
        avg_font = int(np.mean([g["font_size_est"] for g in group]))

        merged.append({
            "text": all_text,
            "bbox": (min_x, min_y, max_x - min_x, max_y - min_y),
            "font_size_est": avg_font,
        })

    return merged


import time

def translate_text_to_french_single(text: str, client) -> str:
    if not text.strip():
        return text

    for attempt in range(3):  # retry up to 3 times
        try:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "Translate the given English text to natural, professional French. "
                            "Return ONLY the French translation, nothing else."
                        ),
                    },
                    {"role": "user", "content": text},
                ],
                temperature=0,
                max_tokens=200,
                timeout = 50,  # 🔥 prevents hanging
            )

            return response.choices[0].message.content.strip()

        except Exception as e:
            print(f"Retry {attempt+1} failed: {e}")
            time.sleep(2)

    return text  # fallback if all retries fail

def translate_image_regions_to_french_batch(texts: list, client) -> list:
    import time as _time
 
    # ← CHANGED: collapse internal newlines so the separator never breaks
    clean_texts = [t.replace("\n", " ").strip() if t else "" for t in texts]
 
    if not any(t.strip() for t in clean_texts):
        return clean_texts
 
    joined = "\n---REGION-END---\n".join(clean_texts)
 
    for attempt in range(3):
        try:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "Translate each English OCR text region into natural French. "
                            "Do not skip short labels, headings, all-caps words, or single words. "
                            "Return ONLY the translated regions, one per line group. "
                            "Keep every region separated exactly with ---REGION-END---. "
                            "Never add extra ---REGION-END--- separators."
                        ),
                    },
                    {"role": "user", "content": joined},
                ],
                temperature=0,
                timeout=30,
            )
 
            translated = response.choices[0].message.content or joined
            parts = [p.strip() for p in translated.split("---REGION-END---")]
 
            if len(parts) == len(clean_texts):
                return parts
 
            print(f"Image batch mismatch: expected {len(clean_texts)}, got {len(parts)}. Falling back to per-region.")
            # ← CHANGED: fall back per-region instead of returning originals
            break
 
        except Exception as e:
            print(f"Image batch translation retry {attempt + 1} failed: {e}")
            _time.sleep(1)
 
    # Per-region fallback
    results = []
    for text in clean_texts:
        if not text.strip():
            results.append(text)
            continue
        try:
            r = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Translate to French. Return only the translation."},
                    {"role": "user", "content": text},
                ],
                temperature=0,
                timeout=15,
            )
            results.append(r.choices[0].message.content.strip())
        except Exception:
            results.append(text)
    return results


# ---------------------------------------------------------------------------
# NEW: Background sampling, complexity detection, and contrast helpers
# ---------------------------------------------------------------------------

def _sample_background_color(pil_image, x, y, w, h):
    """
    Sample the dominant background color under a bounding box by examining
    the border pixels of the region (less likely to contain text pixels).
    Returns an (R, G, B) tuple.
    """
    img_w, img_h = pil_image.size
    x0 = max(0, x);  y0 = max(0, y)
    x1 = min(img_w, x + w);  y1 = min(img_h, y + h)
    if x1 <= x0 or y1 <= y0:
        return (255, 255, 255)
    region = pil_image.crop((x0, y0, x1, y1)).convert("RGB")
    rw, rh = region.size
    if rw < 4 or rh < 4:
        return region.getpixel((rw // 2, rh // 2))
    border_pixels = []
    for bx in range(rw):
        border_pixels.append(region.getpixel((bx, 0)))
        border_pixels.append(region.getpixel((bx, rh - 1)))
    for by in range(1, rh - 1):
        border_pixels.append(region.getpixel((0, by)))
        border_pixels.append(region.getpixel((rw - 1, by)))
    rs = sorted(p[0] for p in border_pixels)
    gs = sorted(p[1] for p in border_pixels)
    bs = sorted(p[2] for p in border_pixels)
    mid = len(rs) // 2
    return (rs[mid], gs[mid], bs[mid])


def _is_complex_background(pil_image, x, y, w, h, variance_threshold=800.0):
    """
    Return True if the region under the bounding box looks like a graph,
    chart, chemical structure, or other complex visual content — as opposed
    to plain text sitting on a flat-colour background.

    Method: measure per-channel pixel variance on the inner 60% of the
    region (avoids border artifacts from adjacent elements).
      High variance  → complex visual content  → skip overlay
      Low variance   → flat background with text → safe to overlay
    """
    img_w, img_h = pil_image.size
    x0 = max(0, x);  y0 = max(0, y)
    x1 = min(img_w, x + w);  y1 = min(img_h, y + h)
    if x1 <= x0 or y1 <= y0:
        return False
    region = pil_image.crop((x0, y0, x1, y1)).convert("RGB")
    rw, rh = region.size
    cx0 = max(0, int(rw * 0.20));  cy0 = max(0, int(rh * 0.20))
    cx1 = min(rw, int(rw * 0.80)); cy1 = min(rh, int(rh * 0.80))
    inner = region.crop((cx0, cy0, cx1, cy1)) if (cx1 > cx0 and cy1 > cy0) else region
    arr = np.array(inner, dtype=np.float32)
    if arr.size == 0:
        return False
    var = float(np.mean([np.var(arr[:, :, c]) for c in range(3)]))
    return var > variance_threshold


def _pick_text_color(bg_color):
    """
    Return 'black' or 'white' for maximum readability against bg_color,
    using the standard luminance formula.
    """
    r, g, b = bg_color
    return "white" if (0.299 * r + 0.587 * g + 0.114 * b) < 128 else "black"


def create_overlay_image(pil_image: Image.Image, translated_regions: list) -> Image.Image:
    result = pil_image.copy()
    draw = ImageDraw.Draw(result)

    # Hard minimum font size — text must always remain readable.
    # There is intentionally NO maximum cap; the loop starts as large as the
    # bounding box allows and shrinks only until the text fits.
    MIN_FONT_SIZE = 10

    for region in translated_regions:
        x, y, w, h = region["bbox"]
        translated = (region["translated"] or "").strip()
        if not translated:
            continue

        # ── Guard: skip regions sitting on graphs / charts / structures ───────
        # We check the original image (not result) so prior overlays don't
        # falsely raise the variance of nearby regions.
        #if _is_complex_background(pil_image, x, y, w, h):
           # print(f"  [overlay] skipped complex-content region: {repr(translated[:40])}")
           # continue

        # ── Sample background color BEFORE erasing anything ───────────────────
        # Reading from pil_image (original) guarantees the color is the true
        # background of the source text, unaffected by earlier overlays.
        bg_color = _sample_background_color(pil_image, x, y, w, h)
        text_color = _pick_text_color(bg_color)

        # ── Tight padding — cover ONLY the text area, not surrounding content ─
        pad_x = max(2, int(w * 0.04))
        pad_y = max(2, int(h * 0.15))

        ex0 = max(0, x - pad_x)
        ey0 = max(0, y - pad_y)
        ex1 = min(result.width,  x + w + pad_x)
        ey1 = min(result.height, y + h + pad_y)
        erase_w = ex1 - ex0
        erase_h = ey1 - ey0

        # ── Fill erased area with the sampled background color ────────────────
        draw.rectangle([ex0, ey0, ex1, ey1], fill=bg_color)

        # ── Font sizing: no maximum, hard minimum of MIN_FONT_SIZE ───────────
        start_size = max(MIN_FONT_SIZE, int(max(erase_h * 0.95, pil_image.height * 0.045)))
        best_font = None
        best_bbox = None

        for size in range(start_size, MIN_FONT_SIZE - 1, -1):
            font = _load_font(size)
            bbox = draw.textbbox((0, 0), translated, font=font)
            text_w = bbox[2] - bbox[0]
            text_h = bbox[3] - bbox[1]
            if text_w <= erase_w * 0.96 and text_h <= erase_h * 0.90:
                best_font = font
                best_bbox = bbox
                break

        # If the text still doesn't fit at MIN_FONT_SIZE, render at MIN anyway.
        # A slightly overflowing label is far better than invisible text.
        if best_font is None:
            best_font = _load_font(MIN_FONT_SIZE)
            best_bbox = draw.textbbox((0, 0), translated, font=best_font)

        text_w = best_bbox[2] - best_bbox[0]
        text_h = best_bbox[3] - best_bbox[1]

        # Centre text inside the erased region
        draw_x = int(ex0 + (erase_w - text_w) / 2)
        draw_y = int(ey0 + (erase_h - text_h) / 2 - best_bbox[1])
        draw.text((draw_x, draw_y), translated, fill=text_color, font=best_font)

    return result


def translate_and_overlay_text(image_path: str, output_path: str, client=None):
    if client is None:
        raise ValueError("OpenAI client must be provided")
 
    text_regions, pil_image = detect_text_regions(image_path)
 
    # pil_image can be None if cv2.imread failed (bad path / unsupported format)
    if pil_image is None:
        return False, f"Could not read image: {image_path}"
 
    if not text_regions:
        # ← CHANGED: save a copy of the original so the file still lands at output_path
        try:
            pil_image.save(output_path)
        except Exception as e:
            return False, f"No text detected and could not save copy: {e}"
        return False, "No text detected in image"
 
    original_texts = [region["text"] for region in text_regions]
    french_texts = translate_image_regions_to_french_batch(original_texts, client)
 
    translated_regions = []
    for region, french in zip(text_regions, french_texts):
        translated_regions.append({
            "original": region["text"],
            "translated": french,
            "bbox": region["bbox"],
            "font_size_est": region.get("font_size_est", 14),
        })
 
    result_image = create_overlay_image(pil_image, translated_regions)
 
    # ← CHANGED: explicitly close source image before saving result
    pil_image.close()
 
    if result_image.mode in ("RGBA", "P"):
        result_image = result_image.convert("RGB")
 
    out_ext = output_path.rsplit(".", 1)[-1].lower()
    if out_ext in {"jpg", "jpeg"}:
        result_image.save(output_path, format="JPEG", quality=95)
    elif out_ext == "bmp":
        result_image.save(output_path, format="BMP")
    elif out_ext in {"tif", "tiff"}:
        result_image.save(output_path, format="TIFF")
    else:
        result_image.save(output_path, format="PNG")
 
    # ← CHANGED: close result image explicitly
    result_image.close()
 
    return True, f"Translated {len(translated_regions)} text region(s)"
# ---------------- BASIC HELPERS ----------------
def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text(filepath: str) -> str:
    """Extract readable text from txt, docx, or pdf."""
    ext = filepath.rsplit(".", 1)[-1].lower()

    try:
        if ext == "txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        if ext == "docx":
            doc = Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs)

        if ext == "pdf":
            reader = PdfReader(filepath)
            parts = []
            for page in reader.pages:
                parts.append(page.extract_text() or "")
            return "\n".join(parts)

        return ""

    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
        return ""


def normalize_whitespace(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def parse_drug_brand_blocks(text: str):
    """
    Handles:

    A) TOP 200 format:
       Generic
       Brand®
       Clinical Use

    B) Class-based format:
       Class
       Generic
       Brand®
    """
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    drugs = []
    i = 0

    while i < len(lines) - 1:
        if i + 2 < len(lines):
            l1, l2, l3 = lines[i], lines[i + 1], lines[i + 2]

            if ("®" in l2) and ("®" not in l1) and ("®" not in l3):
                drugs.append({"generic": l1, "brand": l2, "extra": l3})
                i += 3
                continue

            if ("®" in l3) and ("®" not in l1) and ("®" not in l2):
                drugs.append({"generic": l2, "brand": l3, "extra": l1})
                i += 3
                continue

        l1, l2 = lines[i], lines[i + 1]
        if "®" in l2:
            drugs.append({"generic": l1, "brand": l2, "extra": ""})
            i += 2
        else:
            i += 1

    return drugs


# ---------------- MCQ NORMALIZATION / PARSING ----------------
def build_mcq_normalization_prompt(raw_text: str) -> str:
    return f"""
You are an expert pharmacy exam editor.

Read the text below and extract ALL multiple-choice questions EXACTLY as written.

CRITICAL RULES:
- If a question includes a case scenario, vignette, or clinical stem BEFORE the question sentence,
  include the FULL vignette as part of the question text.
- Do NOT summarize.
- Do NOT shorten.
- Do NOT rewrite unless needed only to repair broken formatting.
- Preserve original wording as much as possible.

OUTPUT RULES:
- Plain text only
- No markdown
- No bullets
- No commentary
- Use EXACTLY this structure for every question:

1. Full question text ending with a question mark?

A) Option text
B) Option text
C) Option text
D) Option text
Ans: C
Tips: Short explanation sentence.

RULES:
- Number sequentially
- Each question must include A) B) C) D)
- Use exactly "Ans:" and "Tips:"
- Leave one blank line between questions
- If the source has no explanation, create a short helpful explanation
- Fix broken option labels if needed
- If source uses A. / B. / C. / D., convert to A) / B) / C) / D)

TEXT:
{raw_text}
""".strip()


def split_text_for_llm(text: str, max_chars: int = 14000, max_chunks: int = 12):
    """
    Split large exam files into chunks while trying to preserve question boundaries.
    """
    text = normalize_whitespace(text)

    # Prefer splitting on numbered questions
    parts = re.split(r"(?=\n?\s*\d+\.\s)", "\n" + text)
    parts = [p.strip() for p in parts if p.strip()]

    if len(parts) <= 1:
        parts = re.split(r"\n\s*\n", text)
        parts = [p.strip() for p in parts if p.strip()]

    chunks = []
    current = ""

    for part in parts:
        candidate = f"{current}\n\n{part}".strip() if current else part
        if len(candidate) <= max_chars:
            current = candidate
        else:
            if current:
                chunks.append(current)
            current = part

    if current:
        chunks.append(current)

    if not chunks:
        chunks = [text[:max_chars]]

    return chunks[:max_chunks]


def normalize_mcq_output_text(text: str) -> str:
    """
    Clean model output into one strict format.
    """
    text = text.replace("```", "")
    text = normalize_whitespace(text)

    lines = []
    for raw_line in text.split("\n"):
        line = raw_line.strip()

        line = re.sub(r"^([A-Da-d])\.\s+", r"\1) ", line)
        line = re.sub(r"^([A-Da-d])\)\s*", lambda m: f"{m.group(1).upper()}) ", line)
        line = re.sub(r"^(Ans|Answer)\s*[:\-]?\s*([A-Da-d])\b", r"Ans: \2", line, flags=re.I)
        line = re.sub(r"^(Tips|Explanation)\s*[:\-]?\s*", "Tips: ", line, flags=re.I)

        lines.append(line)

    cleaned = "\n".join(lines)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def parse_normalized_mcq_text(normalized_text: str):
    """
    Parse strict normalized MCQ text into structured objects.
    """
    normalized_text = normalize_mcq_output_text(normalized_text)

    pattern = re.compile(
        r"""
        ^\s*(\d+)\.\s*(.+?)\n
        A\)\s*(.+?)\n
        B\)\s*(.+?)\n
        C\)\s*(.+?)\n
        D\)\s*(.+?)
        (?:\nAns:\s*([A-D]))?
        (?:\nTips:\s*(.+?))?
        (?=\n\s*\d+\.|\Z)
        """,
        re.S | re.M | re.X,
    )

    mcqs = []
    for match in pattern.finditer(normalized_text):
        question = match.group(2).strip()
        a_opt = match.group(3).strip()
        b_opt = match.group(4).strip()
        c_opt = match.group(5).strip()
        d_opt = match.group(6).strip()
        answer = (match.group(7) or "").strip().upper()
        explanation = (match.group(8) or "No explanation provided.").strip()

        mcqs.append(
            {
                "question": question,
                "options": [
                    f"A) {a_opt}",
                    f"B) {b_opt}",
                    f"C) {c_opt}",
                    f"D) {d_opt}",
                ],
                "answer": answer or "N/A",
                "explanation": explanation,
            }
        )

    return mcqs


def render_normalized_mcq_text(mcqs) -> str:
    """
    Render structured MCQs back into strict normalized text.
    """
    blocks = []

    for idx, mcq in enumerate(mcqs, start=1):
        options = mcq.get("options", [])
        while len(options) < 4:
            label = chr(65 + len(options))
            options.append(f"{label}) ")

        answer = (mcq.get("answer") or "N/A").strip().upper()
        explanation = (mcq.get("explanation") or "No explanation provided.").strip()

        block = "\n".join(
            [
                f"{idx}. {mcq.get('question', '').strip()}",
                options[0].strip(),
                options[1].strip(),
                options[2].strip(),
                options[3].strip(),
                f"Ans: {answer}",
                f"Tips: {explanation}",
            ]
        )
        blocks.append(block)

    return "\n\n".join(blocks).strip()


def regex_extract_mcqs_fallback(text: str):
    """
    Fallback extractor if model normalization fails.
    More robust than the old simple splitter.
    """
    text = normalize_whitespace(text)

    pattern = re.compile(
        r"""
        (?:
            ^|\n
        )
        \s*(\d+)\.\s*(.+?)
        (?=
            \n[A-Da-d][\)\.]
        )
        \n[Aa][\)\.]\s*(.+?)
        \n[Bb][\)\.]\s*(.+?)
        \n[Cc][\)\.]\s*(.+?)
        \n[Dd][\)\.]\s*(.+?)
        (?:
            \n(?:Ans|Answer)\s*[:\-]?\s*([A-Da-d])
        )?
        (?:
            \n(?:Tips|Explanation)\s*[:\-]?\s*(.+?)
        )?
        (?=\n\s*\d+\.|\Z)
        """,
        re.S | re.M | re.X,
    )

    mcqs = []
    for m in pattern.finditer(text):
        question = m.group(2).strip()
        options = [
            f"A) {m.group(3).strip()}",
            f"B) {m.group(4).strip()}",
            f"C) {m.group(5).strip()}",
            f"D) {m.group(6).strip()}",
        ]
        answer = (m.group(7) or "").upper()
        explanation = (m.group(8) or "No explanation provided.").strip()

        mcqs.append(
            {
                "question": question,
                "options": options,
                "answer": answer or "N/A",
                "explanation": explanation,
            }
        )

    return mcqs


def normalize_mcqs_with_gpt(text: str) -> str:
    """
    Use the same CBT-parser standard to normalize MCQs for all generators.
    """
    chunks = split_text_for_llm(text)
    all_mcqs = []

    for chunk in chunks:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You normalize messy exam content into strict MCQ format.",
                },
                {
                    "role": "user",
                    "content": build_mcq_normalization_prompt(chunk),
                },
            ],
            temperature=0,
        )

        chunk_text = normalize_mcq_output_text(response.choices[0].message.content or "")
        chunk_mcqs = parse_normalized_mcq_text(chunk_text)
        all_mcqs.extend(chunk_mcqs)

    return render_normalized_mcq_text(all_mcqs)


def extract_mcqs_with_cbt_standard(text: str):
    """
    One universal MCQ extraction path for:
    - CBT Parser export
    - MCQ Generator 1
    - MCQ Generator 2
    - MCQ Mobile
    - MCQ Grader
    - future MCQ generators
    """
    try:
        normalized_text = normalize_mcqs_with_gpt(text)
        mcqs = parse_normalized_mcq_text(normalized_text)

        if mcqs:
            return mcqs, normalized_text
    except Exception as e:
        print(f"GPT normalization failed, falling back to regex extractor: {e}")

    mcqs = regex_extract_mcqs_fallback(text)
    normalized_text = render_normalized_mcq_text(mcqs) if mcqs else ""
    return mcqs, normalized_text


# ---------------- DOCX OUTPUT HELPERS ----------------
def write_normalized_mcqs_to_docx(normalized_text: str, output_path: str):
    doc = Document()

    def style(run, bold=False, italic=False, color=None):
        run.font.name = "Times New Roman"
        run.font.size = DocxPt(14)
        run.bold = bold
        run.italic = italic
        if color:
            run.font.color.rgb = color

    for line in normalized_text.split("\n"):
        line = line.rstrip()
        p = doc.add_paragraph()
        p.paragraph_format.space_before = DocxPt(0)
        p.paragraph_format.space_after = DocxPt(0)

        m_q = re.match(r"^(\d+\.)(\s*)(.+)", line)
        m_opt = re.match(r"^([A-D]\))(\s*)(.+)", line)

        if m_q:
            style(p.add_run(m_q.group(1)), bold=True)
            p.add_run(m_q.group(2))
            style(p.add_run(m_q.group(3)))
            continue

        if m_opt:
            p.paragraph_format.left_indent = DocxPt(18)
            style(p.add_run(m_opt.group(1)), bold=True)
            p.add_run(m_opt.group(2))
            style(p.add_run(m_opt.group(3)))
            continue

        if line.startswith("Ans:"):
            style(p.add_run("Ans:"), bold=True, color=DocxRGBColor(200, 0, 0))
            p.add_run(" ")
            style(
                p.add_run(line.replace("Ans:", "").strip()),
                color=DocxRGBColor(200, 0, 0),
            )
            continue

        if line.startswith("Tips:"):
            style(p.add_run("Tips:"), bold=True)
            p.add_run(" ")
            style(p.add_run(line.replace("Tips:", "").strip()))
            continue

        style(p.add_run(line))

    doc.save(output_path)


import base64
from docx import Document
from docx.oxml.ns import qn
from lxml import etree

def convert_docx_to_html_faithful(input_path: str) -> str:
    import base64
    from docx import Document
    from docx.oxml.ns import qn
    from lxml import etree
    from docx.text.run import Run

    doc = Document(input_path)

    # ── Image map ──────────────────────────────────────────────────────────
    image_map = {}
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_part = rel.target_part
                ct = img_part.content_type
                ext = ct.split("/")[-1]
                if ext == "jpeg": ext = "jpg"
                b64 = base64.b64encode(img_part.blob).decode("utf-8")
                image_map[rel.rId] = (f"data:{ct};base64,{b64}", ct)
            except Exception as e:
                print(f"Skipping image rel {rel.rId}: {e}")

    # ── Unit helpers ────────────────────────────────────────────────────────
    def emu_to_px(emu):
        try: return round(int(emu) / 914400 * 96)
        except: return None

    def half_pt_to_px(hp):
        try: return round(int(hp) / 2 * 1.333)
        except: return None

    def twip_to_px(tw):
        try: return round(int(tw) / 20 * 1.333)
        except: return None

    def twip_to_pt(tw):
        try: return round(int(tw) / 20, 1)
        except: return None

    def rgb_hex(val):
        if val and len(val) == 6 and val.upper() != "AUTO":
            return f"#{val.lower()}"
        return None

    # ── Numbering helper ────────────────────────────────────────────────────
    _numbering_cache = {}
    def get_list_info(ppr):
        if ppr is None: return None
        num_pr = ppr.find(qn("w:numPr"))
        if num_pr is None: return None
        num_id_el = num_pr.find(qn("w:numId"))
        ilvl_el   = num_pr.find(qn("w:ilvl"))
        if num_id_el is None: return None
        num_id = num_id_el.get(qn("w:val"), "0")
        ilvl   = int(ilvl_el.get(qn("w:val"), "0")) if ilvl_el is not None else 0

        cache_key = (num_id, ilvl)
        if cache_key in _numbering_cache:
            return _numbering_cache[cache_key]

        is_ordered = False
        try:
            nb = doc.part.numbering_part
            ns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
            nxml = nb._element
            num_el = next((n for n in nxml.findall(f'{{{ns}}}num')
                           if n.get(f'{{{ns}}}numId') == num_id), None)
            if num_el is not None:
                abs_id_el = num_el.find(f'{{{ns}}}abstractNumId')
                if abs_id_el is not None:
                    abs_id = abs_id_el.get(f'{{{ns}}}val')
                    abs_el = next((a for a in nxml.findall(f'{{{ns}}}abstractNum')
                                   if a.get(f'{{{ns}}}abstractNumId') == abs_id), None)
                    if abs_el is not None:
                        lvl_el = next((l for l in abs_el.findall(f'{{{ns}}}lvl')
                                       if l.get(f'{{{ns}}}ilvl') == str(ilvl)), None)
                        if lvl_el is not None:
                            fmt_el = lvl_el.find(f'{{{ns}}}numFmt')
                            if fmt_el is not None:
                                fmt = fmt_el.get(f'{{{ns}}}val', '')
                                is_ordered = fmt in ('decimal','lowerLetter','upperLetter',
                                                     'lowerRoman','upperRoman')
        except Exception:
            pass

        result = (num_id, ilvl, is_ordered)
        _numbering_cache[cache_key] = result
        return result

    # ── Inline image extractor ──────────────────────────────────────────────
    def get_inline_images(el):
        imgs = []
        for drawing in el.findall(".//" + qn("w:drawing")):
            for blip in drawing.findall(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
            ):
                r_embed = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                if r_embed and r_embed in image_map:
                    src, ct = image_map[r_embed]
                    width_style = ""
                    for extent in drawing.findall(
                        ".//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}extent"
                    ):
                        cx = extent.get("cx")
                        if cx:
                            px = emu_to_px(cx)
                            if px: width_style = f' style="width:{px}px;max-width:100%;"'
                    imgs.append(f'<img src="{src}"{width_style} />')
        return imgs

    # ── Run parser ──────────────────────────────────────────────────────────
    def parse_run(run_el, para):
        run = Run(run_el, para)
        text = run.text or ""
        if not text: return ""

        escaped = (text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                       .replace('"','&quot;'))

        rpr = run_el.find(qn("w:rPr"))
        styles = {}
        bold = italic = underline = strike = False
        color = bg_color = font_family = None
        font_size_px = None
        vertical = None

        # Check rPr on the run itself
        if rpr is not None:
            if rpr.find(qn("w:b")) is not None:
                b_el = rpr.find(qn("w:b"))
                if b_el.get(qn("w:val"), "true") not in ("false","0"): bold = True
            if rpr.find(qn("w:i")) is not None:
                i_el = rpr.find(qn("w:i"))
                if i_el.get(qn("w:val"), "true") not in ("false","0"): italic = True
            if rpr.find(qn("w:u")) is not None:
                u_el = rpr.find(qn("w:u"))
                if u_el.get(qn("w:val"), "single") not in ("none","false","0"): underline = True
            if rpr.find(qn("w:strike")) is not None: strike = True
            color_el = rpr.find(qn("w:color"))
            if color_el is not None:
                color = rgb_hex(color_el.get(qn("w:val")))
            hl = rpr.find(qn("w:highlight"))
            if hl is not None:
                hl_map = {"yellow":"#ffff00","green":"#00ff00","cyan":"#00ffff",
                          "magenta":"#ff00ff","blue":"#0000ff","red":"#ff0000",
                          "darkBlue":"#00008b","darkCyan":"#008b8b","darkGreen":"#006400",
                          "darkMagenta":"#8b008b","darkRed":"#8b0000","darkYellow":"#808000",
                          "darkGray":"#a9a9a9","lightGray":"#d3d3d3","black":"#000000"}
                bg_color = hl_map.get(hl.get(qn("w:val"),""))
            sz = rpr.find(qn("w:sz"))
            if sz is not None:
                font_size_px = half_pt_to_px(sz.get(qn("w:val")))
            fonts_el = rpr.find(qn("w:rFonts"))
            if fonts_el is not None:
                font_family = (fonts_el.get(qn("w:ascii")) or
                               fonts_el.get(qn("w:hAnsi")) or
                               fonts_el.get(qn("w:cs")))
            vert = rpr.find(qn("w:vertAlign"))
            if vert is not None: vertical = vert.get(qn("w:val"))

        if bold:           styles["font-weight"] = "bold"
        if italic:         styles["font-style"]  = "italic"
        if underline:      styles["text-decoration"] = "underline"
        if strike:
            td = styles.get("text-decoration","")
            styles["text-decoration"] = (td + " line-through").strip()
        if color:          styles["color"] = color
        if bg_color:       styles["background-color"] = bg_color
        if font_size_px:   styles["font-size"] = f"{font_size_px}px"
        if font_family:    styles["font-family"] = f"'{font_family}',sans-serif"

        style_str = ";".join(f"{k}:{v}" for k,v in styles.items())
        html = f'<span style="{style_str}">{escaped}</span>' if style_str else escaped

        if vertical == "superscript": html = f"<sup>{html}</sup>"
        elif vertical == "subscript": html = f"<sub>{html}</sub>"
        return html

    # ── Paragraph parser ────────────────────────────────────────────────────
    def parse_paragraph(para):
        ppr = para._p.find(qn("w:pPr"))
        style_name = ""
        align = None
        indent_left_px = None
        indent_right_px = None
        indent_firstline_px = None
        space_before_pt = None
        space_after_pt = None
        line_height = None
        shading_color = None

        if ppr is not None:
            ps = ppr.find(qn("w:pStyle"))
            if ps is not None: style_name = ps.get(qn("w:val"), "")

            jc = ppr.find(qn("w:jc"))
            if jc is not None:
                jc_map = {"center":"center","right":"right","both":"justify",
                          "distribute":"justify","left":"left"}
                align = jc_map.get(jc.get(qn("w:val"),""), None)

            ind = ppr.find(qn("w:ind"))
            if ind is not None:
                left = ind.get(qn("w:left"))
                right = ind.get(qn("w:right"))
                first = ind.get(qn("w:firstLine"))
                hanging = ind.get(qn("w:hanging"))
                if left: indent_left_px = twip_to_px(left)
                if right: indent_right_px = twip_to_px(right)
                if first: indent_firstline_px = twip_to_px(first)
                if hanging:
                    hp = twip_to_px(hanging)
                    if hp and indent_left_px:
                        indent_firstline_px = -hp

            spacing = ppr.find(qn("w:spacing"))
            if spacing is not None:
                before = spacing.get(qn("w:before"))
                after  = spacing.get(qn("w:after"))
                line   = spacing.get(qn("w:line"))
                line_rule = spacing.get(qn("w:lineRule"))
                if before: space_before_pt = twip_to_pt(before)
                if after:  space_after_pt  = twip_to_pt(after)
                if line:
                    try:
                        lval = int(line)
                        if line_rule in ("auto", None, ""):
                            line_height = round(lval / 240, 2)
                        elif line_rule == "exact":
                            line_height = f"{twip_to_pt(line)}pt"
                        elif line_rule == "atLeast":
                            line_height = f"{twip_to_pt(line)}pt"
                    except: pass

            shd = ppr.find(qn("w:shd"))
            if shd is not None:
                fill = shd.get(qn("w:fill"))
                shading_color = rgb_hex(fill) if fill else None

        # paragraph-level rPr for default run formatting
        p_rpr = ppr.find(qn("w:rPr")) if ppr is not None else None
        p_default_bold   = p_rpr is not None and p_rpr.find(qn("w:b")) is not None
        p_default_color  = None
        p_default_size   = None
        p_default_font   = None
        if p_rpr is not None:
            c_el = p_rpr.find(qn("w:color"))
            if c_el is not None: p_default_color = rgb_hex(c_el.get(qn("w:val")))
            s_el = p_rpr.find(qn("w:sz"))
            if s_el is not None: p_default_size = half_pt_to_px(s_el.get(qn("w:val")))
            f_el = p_rpr.find(qn("w:rFonts"))
            if f_el is not None:
                p_default_font = (f_el.get(qn("w:ascii")) or
                                  f_el.get(qn("w:hAnsi")) or
                                  f_el.get(qn("w:cs")))

        # Build inner html
        inner = ""
        list_info = get_list_info(ppr)

        for child in para._p:
            if callable(child.tag): continue
            local = etree.QName(child).localname
            if local == "r":
                imgs = get_inline_images(child)
                inner += "".join(imgs)
                inner += parse_run(child, para)
            elif local == "hyperlink":
                href = ""
                for attr in ["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", "r:id"]:
                    rel_id = child.get(attr)
                    if rel_id and rel_id in doc.part.rels:
                        href = doc.part.rels[rel_id].target_ref
                        break
                link_inner = ""
                for r_el in child.findall(qn("w:r")):
                    link_inner += parse_run(r_el, para)
                inner += f'<a href="{href}" target="_blank" style="color:#1a0dab;">{link_inner}</a>'
            elif local == "drawing":
                imgs = get_inline_images(child)
                inner += "".join(imgs)
            elif local == "br":
                br_type = child.get(qn("w:type"), "textWrapping")
                if br_type == "page":
                    inner += '<div style="page-break-after:always;"></div>'
                else:
                    inner += "<br>"

        inner = inner.strip()
        if not inner:
            after_pt = space_after_pt if space_after_pt is not None else 0
            return f'<p style="margin:0 0 {after_pt}pt 0;min-height:0.8em;">&nbsp;</p>'

        heading_map = {
            "Heading1":"h1","heading1":"h1","Heading 1":"h1",
            "Heading2":"h2","heading2":"h2","Heading 2":"h2",
            "Heading3":"h3","heading3":"h3","Heading 3":"h3",
            "Heading4":"h4","Heading5":"h5","Heading6":"h6",
            "Title":"h1","Subtitle":"h2",
        }
        html_tag = heading_map.get(style_name, "p")

        p_styles = {}
        if align:                 p_styles["text-align"] = align
        if indent_left_px:        p_styles["padding-left"] = f"{indent_left_px}px"
        if indent_right_px:       p_styles["padding-right"] = f"{indent_right_px}px"
        if indent_firstline_px is not None:
            p_styles["text-indent"] = f"{indent_firstline_px}px"
        if space_before_pt:       p_styles["margin-top"]    = f"{space_before_pt}pt"
        p_styles["margin-bottom"] = f"{space_after_pt}pt" if space_after_pt is not None else "6pt"
        if line_height:
            if isinstance(line_height, float):
                p_styles["line-height"] = str(line_height)
            else:
                p_styles["line-height"] = line_height
        if shading_color and shading_color != "#ffffff":
            p_styles["background-color"] = shading_color
        if p_default_color: p_styles["color"] = p_default_color
        if p_default_size:  p_styles["font-size"] = f"{p_default_size}px"
        if p_default_font:  p_styles["font-family"] = f"'{p_default_font}',sans-serif"

        # List item: use <li> but don't wrap in ul/ol here — handled by flush logic
        if list_info:
            num_id, ilvl, is_ordered = list_info
            indent_px = (ilvl + 1) * 28
            p_styles["margin-left"] = f"{indent_px}px"
            p_styles["padding-left"] = "4px"
            # Remove margin-bottom tightness for lists
            if "margin-bottom" not in p_styles or p_styles["margin-bottom"] == "6pt":
                p_styles["margin-bottom"] = "2pt"
            style_str = ";".join(f"{k}:{v}" for k,v in p_styles.items())
            return f'<li data-numid="{num_id}" data-ilvl="{ilvl}" data-ordered="{1 if is_ordered else 0}" style="{style_str}">{inner}</li>'

        style_str = ";".join(f"{k}:{v}" for k,v in p_styles.items())
        return f'<{html_tag} style="{style_str}">{inner}</{html_tag}>'

    # ── Table parser ────────────────────────────────────────────────────────
    def parse_table(table):
        # Read table-level borders
        tbl_pr = table._tbl.find(qn("w:tblPr"))
        tbl_border_color = "#cccccc"
        tbl_border_size  = 1
        if tbl_pr is not None:
            tb = tbl_pr.find(qn("w:tblBorders"))
            if tb is not None:
                inside_h = tb.find(qn("w:insideH"))
                if inside_h is not None:
                    sz = inside_h.get(qn("w:sz"))
                    col = inside_h.get(qn("w:color"))
                    if sz: tbl_border_size = max(1, round(int(sz) / 8))
                    if col and col.upper() != "AUTO": tbl_border_color = f"#{col.lower()}"

        rows_html = ""
        for row in table.rows:
            cells_html = ""
            for cell in row.cells:
                tc = cell._tc
                # Merge info
                tc_pr = tc.find(qn("w:tcPr"))
                colspan = ""
                cell_width = ""
                cell_bg = ""
                cell_valign = "top"
                cell_border = ""

                if tc_pr is not None:
                    gs = tc_pr.find(qn("w:gridSpan"))
                    if gs is not None:
                        sv = gs.get(qn("w:val"), "1")
                        if sv != "1": colspan = f' colspan="{sv}"'

                    shd = tc_pr.find(qn("w:shd"))
                    if shd is not None:
                        fill = shd.get(qn("w:fill"))
                        if fill and fill.upper() not in ("AUTO","FFFFFF","auto"):
                            cell_bg = f"background-color:#{fill.lower()};"

                    va = tc_pr.find(qn("w:vAlign"))
                    if va is not None:
                        va_map = {"top":"top","center":"middle","bottom":"bottom"}
                        cell_valign = va_map.get(va.get(qn("w:val"),"top"), "top")

                    tcW = tc_pr.find(qn("w:tcW"))
                    if tcW is not None:
                        w_val  = tcW.get(qn("w:w"))
                        w_type = tcW.get(qn("w:type"), "")
                        if w_val and w_type in ("dxa", ""):
                            px = twip_to_px(w_val)
                            if px: cell_width = f"width:{px}px;"

                    # Per-cell borders
                    tc_borders = tc_pr.find(qn("w:tcBorders"))
                    if tc_borders is not None:
                        sides = {}
                        for side in ("top","left","bottom","right"):
                            s_el = tc_borders.find(qn(f"w:{side}"))
                            if s_el is not None:
                                s_val = s_el.get(qn("w:val"),"single")
                                s_color = s_el.get(qn("w:color"),"auto")
                                s_sz    = s_el.get(qn("w:sz"),"8")
                                if s_val == "nil":
                                    sides[side] = "none"
                                else:
                                    px = max(1, round(int(s_sz)/8)) if s_sz.isdigit() else 1
                                    c  = f"#{s_color.lower()}" if s_color.upper() != "AUTO" else tbl_border_color
                                    sides[side] = f"{px}px solid {c}"
                        if sides:
                            for side, val in sides.items():
                                cell_border += f"border-{side}:{val};"

                if not cell_border:
                    cell_border = f"border:{tbl_border_size}px solid {tbl_border_color};"
                cell_inner = "".join(parse_paragraph(p) for p in cell.paragraphs)
                # Nested tables
                for nested_tbl in tc.findall(qn("w:tbl")):
                    from docx.table import Table as DocxTable
                    cell_inner += parse_table(DocxTable(nested_tbl, tc))

                cells_html += (
                    f'<td{colspan} style="padding:5px 8px;vertical-align:{cell_valign};'
                    f'{cell_border}{cell_bg}{cell_width}">{cell_inner}</td>'
                )
            rows_html += f"<tr>{cells_html}</tr>"

        return (
            f'<table style="border-collapse:collapse;width:100%;margin:8px 0;'
            f'table-layout:auto;">{rows_html}</table>'
        )

    # ── Body pass with list flush ───────────────────────────────────────────
    body_html = ""
    body_el   = doc.element.body
    para_idx  = 0
    table_idx = 0

    # list state
    list_stack = []   # stack of (num_id, ilvl, is_ordered)

    def flush_lists_to(target_ilvl=-1):
        nonlocal body_html
        while list_stack and list_stack[-1][1] > target_ilvl:
            _, ilvl, is_ordered = list_stack.pop()
            tag = "ol" if is_ordered else "ul"
            body_html += f"</{tag}>"

    for child in body_el:
        if callable(child.tag): continue
        try:
            local = etree.QName(child).localname
        except Exception:
            continue

        if local == "p":
            if para_idx >= len(doc.paragraphs):
                continue
            try:
                para      = doc.paragraphs[para_idx]
                para_idx += 1
                ppr       = para._p.find(qn("w:pPr"))
                li = get_list_info(ppr)

                if li:
                    num_id, ilvl, is_ordered = li
                    while list_stack and list_stack[-1][1] >= ilvl:
                        _, _, ordered = list_stack.pop()
                        tag = "ol" if ordered else "ul"
                        body_html += f"</{tag}>"
                    tag = "ol" if is_ordered else "ul"
                    body_html += f'<{tag} style="margin:4pt 0 4pt {ilvl*28}px;padding-left:20px;">'
                    list_stack.append((num_id, ilvl, is_ordered))
                    body_html += parse_paragraph(para)
                else:
                    flush_lists_to(-1)
                    body_html += parse_paragraph(para)
            except Exception as e:
                print(f"Paragraph parse error (idx {para_idx}): {e}")
                para_idx += 1

        elif local == "tbl":
            flush_lists_to(-1)
            if table_idx < len(doc.tables):
                try:
                    body_html += parse_table(doc.tables[table_idx])
                except Exception as e:
                    print(f"Table parse error (idx {table_idx}): {e}")
            table_idx += 1

        elif local == "sectPr":
            pass

    flush_lists_to(-1)

    # ── Document-level page layout ──────────────────────────────────────────
    page_width_px  = 816   # default Letter
    margin_left_px = 96
    margin_right_px = 96
    try:
        body_pr = body_el.find(qn("w:sectPr"))
        if body_pr is None:
            # last child may be sectPr
            body_pr = body_el[-1] if len(body_el) > 0 else None
            if body_pr is not None and etree.QName(body_pr).localname != "sectPr":
                body_pr = None
        if body_pr is not None:
            pg_sz = body_pr.find(qn("w:pgSz"))
            pg_mar = body_pr.find(qn("w:pgMar"))
            if pg_sz is not None:
                w = pg_sz.get(qn("w:w"))
                if w: page_width_px = twip_to_px(w) or page_width_px
            if pg_mar is not None:
                ml = pg_mar.get(qn("w:left"))
                mr = pg_mar.get(qn("w:right"))
                if ml: margin_left_px  = twip_to_px(ml) or margin_left_px
                if mr: margin_right_px = twip_to_px(mr) or margin_right_px
    except Exception as e:
        print(f"Page layout read failed: {e}")

    content_width_px = page_width_px - margin_left_px - margin_right_px

    # ── Collect default body font/size from document styles ─────────────────
    default_font = "Calibri, Arial, sans-serif"
    default_size = "12pt"
    try:
        from docx.oxml.ns import qn as _qn
        styles_part = doc.part.styles
        if styles_part is not None:
            ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            for style_el in styles_part._element.findall(f"{{{ns}}}style"):
                if (style_el.get(f"{{{ns}}}type") == "paragraph" and
                        style_el.get(f"{{{ns}}}default") == "1"):
                    rpr = style_el.find(f"{{{ns}}}rPr")
                    if rpr is not None:
                        sz = rpr.find(f"{{{ns}}}sz")
                        fn = rpr.find(f"{{{ns}}}rFonts")
                        if sz is not None:
                            pt = round(int(sz.get(f"{{{ns}}}val","24")) / 2, 1)
                            default_size = f"{pt}pt"
                        if fn is not None:
                            fname = (fn.get(f"{{{ns}}}ascii") or
                                     fn.get(f"{{{ns}}}hAnsi") or
                                     fn.get(f"{{{ns}}}cs"))
                            if fname: default_font = f"'{fname}', sans-serif"
                    break
    except Exception:
        pass

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1.0"/>
<title>Document</title>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    font-family: {default_font};
    font-size: {default_size};
    line-height: 1.5;
    color: #111;
    background: #fff;
    padding: {margin_left_px}px;
    max-width: {page_width_px}px;
    margin: 0 auto;
  }}
  p, h1, h2, h3, h4, h5, h6 {{ margin: 0; padding: 0; }}
  table {{ border-collapse: collapse; }}
  img {{ display: inline-block; vertical-align: middle; max-width: 100%; }}
  a {{ color: #1a0dab; }}
  ol, ul {{ margin: 0; padding: 0; }}
  li {{ margin: 0; }}
</style>
</head>
<body>
{body_html}
</body>
</html>"""

    return html


def convert_docx_to_html_libreoffice(input_path: str, output_path: str) -> str:
    """
    Pixel-perfect DOCX → HTML using LibreOffice Writer macro export.
    Uses the 'writer8' + fixed-layout HTML approach that preserves
    absolute positioning of every element exactly like wordtohtml.net
    """
    import subprocess, shutil, tempfile, glob, base64, re as _re
    from pathlib import Path

    def find_bin(*names):
        for n in names:
            found = shutil.which(n)
            if found:
                return found
            if Path(n).exists():
                return n
        return None

    lo_bin = find_bin(
        "libreoffice", "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/libreoffice", "/usr/bin/soffice",
        "/usr/lib/libreoffice/program/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    )

    if lo_bin:
        with tempfile.TemporaryDirectory() as tmpdir:
            try:
                # Use the fixed-layout HTML filter — this is what gives
                # pixel-perfect absolute positioning like wordtohtml.net
                filter_options = (
                    "HTML (StarWriter):"
                    "EmbedImages=true,"
                    "FixedLayout=true,"
                    "WriterLayout=true"
                )
                result = subprocess.run([
                    lo_bin,
                    "--headless",
                    "--norestore",
                    "--nofirststartwizard",
                    "--convert-to", f"html:{filter_options}",
                    "--outdir", tmpdir,
                    input_path,
                ], capture_output=True, text=True, timeout=120)

                print("LO stdout:", result.stdout)
                print("LO stderr:", result.stderr)

                html_files = glob.glob(os.path.join(tmpdir, "*.html")) + \
                             glob.glob(os.path.join(tmpdir, "*.HTML"))

                if not html_files:
                    # Retry with simpler filter name
                    result2 = subprocess.run([
                        lo_bin,
                        "--headless",
                        "--norestore",
                        "--convert-to", "html",
                        "--outdir", tmpdir,
                        input_path,
                    ], capture_output=True, text=True, timeout=120)
                    html_files = glob.glob(os.path.join(tmpdir, "*.html")) + \
                                 glob.glob(os.path.join(tmpdir, "*.HTML"))

                if html_files:
                    html_file = html_files[0]
                    with open(html_file, "r", encoding="utf-8", errors="replace") as f:
                        raw_html = f.read()

                    # ── Embed ALL images as base64 ─────────────────────────
                    # LO puts images in a subfolder e.g. filename_html_files/
                    base_name = Path(input_path).stem
                    img_dirs = [
                        tmpdir,
                        os.path.dirname(html_file),
                        os.path.join(tmpdir, base_name + "_html_files"),
                    ]
                    img_dirs += glob.glob(os.path.join(tmpdir, "*_html_files"))
                    img_dirs += glob.glob(os.path.join(tmpdir, "*_files"))
                    img_dirs = list(dict.fromkeys(img_dirs))  # dedupe

                    def embed_src(match):
                        src = match.group(1)
                        if src.startswith(("http", "data:", "//", "#")):
                            return match.group(0)
                        decoded = src
                        for d in img_dirs:
                            for candidate in [
                                os.path.join(d, decoded),
                                os.path.join(d, os.path.basename(decoded)),
                            ]:
                                if os.path.exists(candidate):
                                    ext = candidate.rsplit(".", 1)[-1].lower()
                                    mime = {
                                        "png": "image/png",
                                        "jpg": "image/jpeg",
                                        "jpeg": "image/jpeg",
                                        "gif": "image/gif",
                                        "svg": "image/svg+xml",
                                        "bmp": "image/bmp",
                                        "tiff": "image/tiff",
                                        "webp": "image/webp",
                                        "wmf": "image/x-wmf",
                                        "emf": "image/x-emf",
                                    }.get(ext, "image/png")
                                    with open(candidate, "rb") as imgf:
                                        b64 = base64.b64encode(imgf.read()).decode()
                                    return f'src="data:{mime};base64,{b64}"'
                        return match.group(0)

                    raw_html = _re.sub(r'src="([^"]+)"', embed_src, raw_html)

                    # ── Also embed any url(...) in <style> blocks ──────────
                    def embed_css_url(match):
                        src = match.group(1).strip("'\"")
                        if src.startswith(("http", "data:", "//")):
                            return match.group(0)
                        for d in img_dirs:
                            for candidate in [
                                os.path.join(d, src),
                                os.path.join(d, os.path.basename(src)),
                            ]:
                                if os.path.exists(candidate):
                                    ext = candidate.rsplit(".", 1)[-1].lower()
                                    mime = {
                                        "png": "image/png", "jpg": "image/jpeg",
                                        "jpeg": "image/jpeg", "gif": "image/gif",
                                        "svg": "image/svg+xml", "bmp": "image/bmp",
                                        "woff": "font/woff", "woff2": "font/woff2",
                                        "ttf": "font/ttf",
                                    }.get(ext, "image/png")
                                    with open(candidate, "rb") as imgf:
                                        b64 = base64.b64encode(imgf.read()).decode()
                                    return f"url('data:{mime};base64,{b64}')"
                        return match.group(0)

                    raw_html = _re.sub(r'url\(([^)]+)\)', embed_css_url, raw_html)

                    # ── Inject responsive wrapper + scrollbar fix ──────────
                    inject_css = """
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<style>
  html { background: #f0f0f0; }
  body {
    margin: 0 auto !important;
    padding: 20px !important;
    background: #ffffff !important;
    box-shadow: 0 0 20px rgba(0,0,0,0.15);
    overflow-x: auto !important;
  }
  img { max-width: 100% !important; height: auto !important; }
  table { border-collapse: collapse; max-width: 100%; }
  @media print {
    html { background: white; }
    body { box-shadow: none; padding: 0 !important; margin: 0 !important; }
  }
</style>"""
                    if "</head>" in raw_html:
                        raw_html = raw_html.replace("</head>", inject_css + "\n</head>", 1)
                    else:
                        raw_html = inject_css + raw_html

                    with open(output_path, "w", encoding="utf-8") as out:
                        out.write(raw_html)

                    return raw_html

            except subprocess.TimeoutExpired:
                print("LibreOffice timed out")
            except Exception as e:
                import traceback
                print(f"LibreOffice conversion failed: {e}")
                print(traceback.format_exc())

    # ── Fallback: manual parser ────────────────────────────────────────────
    print("LibreOffice not found or failed — using manual parser fallback")
    html_output = convert_docx_to_html_faithful(input_path)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_output)
    return html_output
# ---------------- PPT HELPERS ----------------
def set_text_frame_defaults(text_frame, margins=(20, 20, 18, 18), align=PP_ALIGN.LEFT):
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(margins[0])
    text_frame.margin_right = Pt(margins[1])
    text_frame.margin_top = Pt(margins[2])
    text_frame.margin_bottom = Pt(margins[3])
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if text_frame.paragraphs:
        text_frame.paragraphs[0].alignment = align


def question_font_size(question: str, large=34, medium=30, small=26, tiny=22):
    n = len(re.sub(r"\s+", " ", question).strip())
    if n > 520:
        return Pt(tiny)
    if n > 360:
        return Pt(small)
    if n > 220:
        return Pt(medium)
    return Pt(large)


def option_font_size(option_text: str, large=24, medium=21, small=18):
    n = len(option_text.strip())
    if n > 80:
        return Pt(small)
    if n > 50:
        return Pt(medium)
    return Pt(large)


def ensure_four_options(mcq: dict):
    options = list(mcq.get("options", []))
    while len(options) < 4:
        label = chr(65 + len(options))
        options.append(f"{label}) ")
    return options[:4]


def shrink_text_to_fit(text_frame, paragraph, min_size=22):
    while paragraph.font.size and paragraph.font.size.pt > min_size:
        if len(text_frame.text.split("\n")) <= 4:
            break
        paragraph.font.size = Pt(paragraph.font.size.pt - 2)


# ---------------- PPT TEMPLATES ----------------
def create_vba_template_presentation(mcqs, output_path):
    if not PPTX_AVAILABLE:
        return False

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        BLUE = PptRGBColor(100, 149, 237)
        CUSTOMBG = PptRGBColor(249, 248, 242)
        DARK = PptRGBColor(22, 22, 26)
        WHITE = PptRGBColor(255, 255, 255)

        for i, mcq in enumerate(mcqs, 1):
            slide1 = prs.slides.add_slide(prs.slide_layouts[6])

            bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg1.fill.solid()
            bg1.fill.fore_color.rgb = CUSTOMBG
            bg1.line.fill.background()

            qBox1 = slide1.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.7), Inches(12.3), Inches(2.0)
            )
            qBox1.fill.solid()
            qBox1.fill.fore_color.rgb = CUSTOMBG
            qBox1.line.color.rgb = BLUE
            qBox1.line.width = Pt(8)
            qBox1.adjustments[0] = 0.08

            qText1 = qBox1.text_frame
            qText1.clear()
            qText1.word_wrap = True
            qText1.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            qText1.margin_left = Pt(25)
            qText1.margin_right = Pt(25)
            qText1.margin_top = Pt(20)
            qText1.margin_bottom = Pt(20)

            p = qText1.paragraphs[0]
            p.text = f"Q{i}: {mcq['question']}"
            p.font.name = "Arial"
            p.font.size = question_font_size(mcq["question"], large=36, medium=32, small=28, tiny=24)
            p.font.bold = True
            p.font.color.rgb = DARK

            badge1 = slide1.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.15), Inches(0.45), Inches(1.1), Inches(1.1)
            )
            badge1.fill.solid()
            badge1.fill.fore_color.rgb = BLUE
            badge1.line.fill.background()

            badgeText1 = badge1.text_frame
            badgeText1.text = str(i)
            badgeText1.paragraphs[0].font.name = "Arial"
            badgeText1.paragraphs[0].font.size = Pt(36)
            badgeText1.paragraphs[0].font.bold = True
            badgeText1.paragraphs[0].font.color.rgb = WHITE
            badgeText1.paragraphs[0].alignment = PP_ALIGN.CENTER

            option_top = Inches(3.0)
            options = ensure_four_options(mcq)

            for j, option in enumerate(options):
                option_box = slide1.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.5),
                    option_top + (j * Inches(1.0)),
                    Inches(12.3),
                    Inches(0.85),
                )
                option_box.fill.solid()
                option_box.fill.fore_color.rgb = CUSTOMBG
                option_box.line.color.rgb = BLUE
                option_box.line.width = Pt(8)
                option_box.adjustments[0] = 0.08

                option_text = option_box.text_frame
                option_text.word_wrap = True
                option_text.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                option_text.margin_left = Pt(25)
                option_text.margin_right = Pt(25)
                option_text.margin_top = Pt(12)
                option_text.margin_bottom = Pt(12)
                option_text.text = option

                option_text.paragraphs[0].font.name = "Arial"
                option_text.paragraphs[0].font.size = option_font_size(option)
                option_text.paragraphs[0].font.italic = True
                option_text.paragraphs[0].font.color.rgb = DARK

            slide2 = prs.slides.add_slide(prs.slide_layouts[6])

            bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg2.fill.solid()
            bg2.fill.fore_color.rgb = CUSTOMBG
            bg2.line.fill.background()

            qBox2 = slide2.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.7), Inches(12.3), Inches(2.0)
            )
            qBox2.fill.solid()
            qBox2.fill.fore_color.rgb = CUSTOMBG
            qBox2.line.color.rgb = BLUE
            qBox2.line.width = Pt(8)
            qBox2.adjustments[0] = 0.08

            qText2 = qBox2.text_frame
            qText2.clear()
            qText2.word_wrap = True
            qText2.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            qText2.margin_left = Pt(25)
            qText2.margin_right = Pt(25)
            qText2.margin_top = Pt(20)
            qText2.margin_bottom = Pt(20)

            p = qText2.paragraphs[0]
            p.text = f"Q{i}: {mcq['question']}"
            p.font.name = "Arial"
            p.font.size = question_font_size(mcq["question"], large=36, medium=32, small=28, tiny=24)
            p.font.bold = True
            p.font.color.rgb = DARK

            badge2 = slide2.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.15), Inches(0.45), Inches(1.1), Inches(1.1)
            )
            badge2.fill.solid()
            badge2.fill.fore_color.rgb = BLUE
            badge2.line.fill.background()

            badgeText2 = badge2.text_frame
            badgeText2.text = str(i)
            badgeText2.paragraphs[0].font.name = "Arial"
            badgeText2.paragraphs[0].font.size = Pt(36)
            badgeText2.paragraphs[0].font.bold = True
            badgeText2.paragraphs[0].font.color.rgb = WHITE
            badgeText2.paragraphs[0].alignment = PP_ALIGN.CENTER

            answer_box = slide2.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5),
                Inches(3.0),
                Inches(12.3),
                Inches(3.2),
            )
            answer_box.fill.solid()
            answer_box.fill.fore_color.rgb = CUSTOMBG
            answer_box.line.color.rgb = BLUE
            answer_box.line.width = Pt(8)
            answer_box.adjustments[0] = 0.08

            answer_content = []
            if mcq["answer"]:
                answer_content.append(f"Correct Answer: {mcq['answer']}")
            if mcq["explanation"]:
                answer_content.append(f"Explanation: {mcq['explanation']}")

            answer_text = answer_box.text_frame
            answer_text.text = "\n\n".join(answer_content)
            answer_text.word_wrap = True
            answer_text.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            answer_text.margin_left = Pt(30)
            answer_text.margin_right = Pt(30)
            answer_text.margin_top = Pt(25)
            answer_text.margin_bottom = Pt(25)

            for paragraph in answer_text.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(28)
                paragraph.font.color.rgb = DARK
                paragraph.alignment = PP_ALIGN.LEFT

            if answer_text.paragraphs:
                answer_text.paragraphs[0].font.bold = True
                answer_text.paragraphs[0].font.italic = True

        prs.save(output_path)
        return True

    except Exception as e:
        print(f"Error creating VBA template presentation: {e}")
        return False


def create_mcq_generator2_exact(mcqs, output_path):
    """
    MCQ Generator 2
    - balanced blue question slide
    - clean red answer slide
    - adaptive text sizing
    - better spacing and proportions
    """
    if not PPTX_AVAILABLE:
        return False

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        BG = PptRGBColor(239, 243, 248)
        RIGHT_BG = PptRGBColor(240, 243, 246)
        BLUE = PptRGBColor(45, 105, 220)
        RED = PptRGBColor(210, 35, 35)
        WHITE = PptRGBColor(255, 255, 255)
        DARK = PptRGBColor(24, 28, 35)
        CARD_FILL = PptRGBColor(255, 255, 255)
        CARD_LINE = PptRGBColor(200, 207, 216)

        left_panel_width = Inches(6.35)
        right_start = left_panel_width
        right_card_x = Inches(7.35)
        right_card_w = Inches(5.15)

        def fit_question_size(text):
            n = len((text or "").strip())
            if n > 500:
                return Pt(18)
            if n > 350:
                return Pt(22)
            if n > 220:
                return Pt(25)
            return Pt(29)

        def fit_option_size(text):
            n = len((text or "").strip())
            if n > 90:
                return Pt(15)
            if n > 65:
                return Pt(17)
            return Pt(20)

        def fit_explanation_size(text):
            n = len((text or "").strip())
            if n > 700:
                return Pt(15)
            if n > 500:
                return Pt(17)
            if n > 320:
                return Pt(19)
            return Pt(21)

        for i, mcq in enumerate(mcqs, 1):
            options = ensure_four_options(mcq)

            # =========================
            # SLIDE 1 — QUESTION
            # =========================
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = BG
            bg.line.fill.background()

            left_panel = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                left_panel_width,
                prs.slide_height,
            )
            left_panel.fill.solid()
            left_panel.fill.fore_color.rgb = BLUE
            left_panel.line.fill.background()

            right_panel = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                right_start,
                0,
                prs.slide_width - left_panel_width,
                prs.slide_height,
            )
            right_panel.fill.solid()
            right_panel.fill.fore_color.rgb = RIGHT_BG
            right_panel.line.fill.background()

            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.55),
                Inches(0.50),
                Inches(0.78),
                Inches(0.78),
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = WHITE
            num_circle.line.fill.background()

            num_tf = num_circle.text_frame
            num_tf.clear()
            num_tf.text = str(i)
            num_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_tf.paragraphs[0].font.name = "Arial"
            num_tf.paragraphs[0].font.size = Pt(24)
            num_tf.paragraphs[0].font.bold = True
            num_tf.paragraphs[0].font.color.rgb = BLUE

            q_box = slide.shapes.add_textbox(
                Inches(0.78),
                Inches(1.00),
                Inches(5.00),
                Inches(5.95),
            )
            q_tf = q_box.text_frame
            q_tf.clear()
            q_tf.word_wrap = True
            q_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            q_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            q_tf.margin_left = Pt(20)
            q_tf.margin_right = Pt(20)
            q_tf.margin_top = Pt(20)
            q_tf.margin_bottom = Pt(20)

            q_p = q_tf.paragraphs[0]
            q_p.text = f"{i}. {mcq['question']}"
            q_p.alignment = PP_ALIGN.LEFT
            q_p.font.name = "Arial"
            q_p.font.size = fit_question_size(mcq["question"])
            q_p.font.bold = True
            q_p.font.color.rgb = WHITE

            option_y_positions = [Inches(0.90), Inches(2.08), Inches(3.26), Inches(4.44)]

            for idx, option in enumerate(options):
                opt_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    right_card_x,
                    option_y_positions[idx],
                    right_card_w,
                    Inches(0.92),
                )
                opt_box.fill.solid()
                opt_box.fill.fore_color.rgb = CARD_FILL
                opt_box.line.color.rgb = CARD_LINE
                opt_box.line.width = Pt(1.3)
                opt_box.adjustments[0] = 0.06

                opt_tf = opt_box.text_frame
                opt_tf.clear()
                opt_tf.word_wrap = True
                opt_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                opt_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                opt_tf.margin_left = Pt(18)
                opt_tf.margin_right = Pt(18)
                opt_tf.margin_top = Pt(8)
                opt_tf.margin_bottom = Pt(8)

                opt_p = opt_tf.paragraphs[0]
                opt_p.text = option
                opt_p.alignment = PP_ALIGN.CENTER
                opt_p.font.name = "Arial"
                opt_p.font.size = fit_option_size(option)
                opt_p.font.bold = True
                opt_p.font.color.rgb = DARK

            # =========================
            # SLIDE 2 — ANSWER
            # =========================
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])

            bg2 = slide2.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            bg2.fill.solid()
            bg2.fill.fore_color.rgb = BG
            bg2.line.fill.background()

            correct_letter = (mcq.get("answer") or "N/A").strip().upper()
            correct_text = next(
                (o[3:].strip() for o in options if o.startswith(correct_letter)),
                "",
            )

            answer_box = slide2.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.75),
                Inches(0.72),
                Inches(7.20),
                Inches(1.35),
            )
            answer_box.fill.solid()
            answer_box.fill.fore_color.rgb = RED
            answer_box.line.fill.background()

            ans_tf = answer_box.text_frame
            ans_tf.clear()
            ans_tf.word_wrap = True
            ans_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            ans_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            ans_tf.margin_left = Pt(18)
            ans_tf.margin_right = Pt(18)
            ans_tf.margin_top = Pt(10)
            ans_tf.margin_bottom = Pt(10)

            ans_p = ans_tf.paragraphs[0]
            ans_p.text = f"Answer: {correct_letter}) {correct_text}"
            ans_p.alignment = PP_ALIGN.LEFT
            ans_p.font.name = "Arial"
            ans_p.font.size = Pt(24)
            ans_p.font.bold = True
            ans_p.font.color.rgb = WHITE

            exp_box = slide2.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.75),
                Inches(2.40),
                Inches(11.85),
                Inches(2.45),
            )
            exp_box.fill.solid()
            exp_box.fill.fore_color.rgb = WHITE
            exp_box.line.color.rgb = CARD_LINE
            exp_box.line.width = Pt(1.3)
            exp_box.adjustments[0] = 0.06

            exp_tf = exp_box.text_frame
            exp_tf.clear()
            exp_tf.word_wrap = True
            exp_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            exp_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            exp_tf.margin_left = Pt(24)
            exp_tf.margin_right = Pt(24)
            exp_tf.margin_top = Pt(18)
            exp_tf.margin_bottom = Pt(18)

            exp_title = exp_tf.paragraphs[0]
            exp_title.text = "Explanation:"
            exp_title.alignment = PP_ALIGN.LEFT
            exp_title.font.name = "Arial"
            exp_title.font.size = Pt(20)
            exp_title.font.bold = True
            exp_title.font.color.rgb = DARK

            exp_body = exp_tf.add_paragraph()
            exp_body.text = mcq.get("explanation") or "No explanation provided."
            exp_body.alignment = PP_ALIGN.LEFT
            exp_body.font.name = "Arial"
            exp_body.font.size = fit_explanation_size(mcq.get("explanation") or "")
            exp_body.font.color.rgb = DARK

        prs.save(output_path)
        return True

    except Exception as e:
        print(f"MCQ2 ERROR: {e}")
        return False
def create_mcq_generator3_exact(mcqs, output_path):
    """
    MCQ Generator 3
    - question slide: teal left panel, larger option circles/text
    - answer slide: centered answer row with checkmark beside answer text
    - improved auto-scaling
    """
    print(f"[MCQ3] Starting: {len(mcqs)} MCQs", flush=True)

    if not PPTX_AVAILABLE:
        print("[MCQ3] pptx not available", flush=True)
        return False

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        TEAL = PptRGBColor(18, 138, 150)
        GREEN = PptRGBColor(98, 190, 72)
        WHITE = PptRGBColor(255, 255, 255)
        OFF_WHITE = PptRGBColor(245, 245, 245)
        DARK = PptRGBColor(38, 43, 49)
        MID_GRAY = PptRGBColor(120, 120, 120)

        slide_w = prs.slide_width
        slide_h = prs.slide_height
        left_panel_w = Inches(4.0)

        def clean_text(value):
            value = "" if value is None else str(value)
            value = value.replace("\x00", "")
            value = value.replace("\r\n", "\n").replace("\r", "\n")
            value = " ".join(value.split())
            return value.strip()

        def ensure_answer_letter(ans):
            ans = clean_text(ans).upper()
            return ans if ans in {"A", "B", "C", "D"} else "N/A"

        def strip_option_label(option):
            option = clean_text(option)
            if len(option) >= 3 and option[0].upper() in {"A", "B", "C", "D"} and option[1] == ")":
                return option[3:].strip()
            return option

        def set_text_frame(tf, ml=0, mr=0, mt=0, mb=0, wrap=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE):
            tf.clear()
            tf.word_wrap = wrap
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.vertical_anchor = valign
            tf.margin_left = Pt(ml)
            tf.margin_right = Pt(mr)
            tf.margin_top = Pt(mt)
            tf.margin_bottom = Pt(mb)

        def add_full_bg(slide, color):
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()
            return bg

        def add_circle(slide, left, top, size, fill_rgb, text="", text_rgb=WHITE, font_size=22, bold=True):
            shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill_rgb
            shp.line.fill.background()

            tf = shp.text_frame
            set_text_frame(tf, 0, 0, 0, 0)

            p = tf.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = text_rgb
            return shp

        def fit_question_size(text):
            n = len(clean_text(text))
            if n > 560:
                return Pt(14)
            if n > 480:
                return Pt(16)
            if n > 400:
                return Pt(18)
            if n > 320:
                return Pt(20)
            if n > 250:
                return Pt(22)
            if n > 180:
                return Pt(25)
            return Pt(30)

        def fit_option_size(text):
            n = len(clean_text(text))
            if n > 170:
                return Pt(17)
            if n > 130:
                return Pt(19)
            if n > 95:
                return Pt(21)
            if n > 65:
                return Pt(24)
            return Pt(27)

        def fit_answer_size(text):
            n = len(clean_text(text))
            if n > 220:
                return Pt(18)
            if n > 175:
                return Pt(20)
            if n > 135:
                return Pt(23)
            if n > 95:
                return Pt(27)
            if n > 60:
                return Pt(32)
            return Pt(38)

        def fit_explanation_size(text):
            n = len(clean_text(text))
            if n > 420:
                return Pt(11)
            if n > 320:
                return Pt(12)
            if n > 240:
                return Pt(13)
            return Pt(15)

        for i, mcq in enumerate(mcqs, start=1):
            print(f"[MCQ3] Q{i}", flush=True)

            question = clean_text(mcq.get("question", ""))
            explanation = clean_text(mcq.get("explanation", "")) or "No explanation provided."
            answer_letter = ensure_answer_letter(mcq.get("answer", ""))

            options = [clean_text(x) for x in ensure_four_options(mcq)]
            option_texts = [strip_option_label(x) for x in options]

            correct_text = ""
            if answer_letter in {"A", "B", "C", "D"}:
                idx = ord(answer_letter) - ord("A")
                if 0 <= idx < len(option_texts):
                    correct_text = option_texts[idx]

            # =====================================================
            # QUESTION SLIDE
            # =====================================================
            slide1 = prs.slides.add_slide(prs.slide_layouts[6])
            add_full_bg(slide1, OFF_WHITE)

            left_panel = slide1.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0,
                0,
                left_panel_w,
                slide_h,
            )
            left_panel.fill.solid()
            left_panel.fill.fore_color.rgb = TEAL
            left_panel.line.fill.background()

            add_circle(
                slide1,
                Inches(0.42),
                Inches(0.50),
                Inches(0.62),
                WHITE,
                text=str(i),
                text_rgb=TEAL,
                font_size=20,
            )

            q_box = slide1.shapes.add_textbox(
                Inches(0.65),
                Inches(1.05),
                Inches(3.05),
                Inches(5.95),
            )
            q_tf = q_box.text_frame
            set_text_frame(q_tf, 0, 0, 0, 0, wrap=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)

            q_p = q_tf.paragraphs[0]
            q_p.text = f"{i}. {question}"
            q_p.alignment = PP_ALIGN.LEFT
            q_p.font.name = "Arial"
            q_p.font.size = fit_question_size(question)
            q_p.font.bold = True
            q_p.font.color.rgb = WHITE

            # Bigger option circles + bigger option text
            option_y_positions = [
                Inches(0.95),
                Inches(2.25),
                Inches(3.55),
                Inches(4.85),
            ]

            letter_circle_x = Inches(4.38)
            option_text_x = Inches(5.12)
            option_text_w = Inches(7.35)
            option_circle_size = Inches(0.68)

            for idx, opt_text in enumerate(option_texts):
                label = chr(65 + idx)

                add_circle(
                    slide1,
                    letter_circle_x,
                    option_y_positions[idx] - Inches(0.06),
                    option_circle_size,
                    TEAL,
                    text=label,
                    text_rgb=WHITE,
                    font_size=22,
                )

                tb = slide1.shapes.add_textbox(
                    option_text_x,
                    option_y_positions[idx] - Inches(0.08),
                    option_text_w,
                    Inches(0.90),
                )
                tf = tb.text_frame
                set_text_frame(tf, 0, 0, 0, 0, wrap=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)

                p = tf.paragraphs[0]
                p.text = opt_text
                p.alignment = PP_ALIGN.LEFT
                p.font.name = "Arial"
                p.font.size = fit_option_size(opt_text)
                p.font.bold = True
                p.font.color.rgb = DARK

            # =====================================================
            # ANSWER SLIDE
            # =====================================================
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            add_full_bg(slide2, WHITE)

            answer_text = correct_text if correct_text else f"Correct answer: {answer_letter}"
            answer_font_size = fit_answer_size(answer_text)

            # Centered row: letter circle + answer text + check circle
            row_left = Inches(0.90)
            row_top = Inches(2.30)
            row_width = Inches(11.50)
            row_height = Inches(1.25)

            answer_letter_size = Inches(0.66)
            check_size = Inches(0.66)

            add_circle(
                slide2,
                row_left,
                row_top + Inches(0.24),
                answer_letter_size,
                GREEN,
                text=answer_letter if answer_letter in {"A", "B", "C", "D"} else "",
                text_rgb=WHITE,
                font_size=22,
            )

            answer_box_left = row_left + Inches(0.88)
            answer_box_width = row_width - Inches(1.90)

            answer_box = slide2.shapes.add_textbox(
                answer_box_left,
                row_top,
                answer_box_width,
                row_height,
            )
            ans_tf = answer_box.text_frame
            set_text_frame(ans_tf, 0, 0, 0, 0, wrap=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)

            ans_p = ans_tf.paragraphs[0]
            ans_p.text = answer_text
            ans_p.alignment = PP_ALIGN.CENTER
            ans_p.font.name = "Arial"
            ans_p.font.size = answer_font_size
            ans_p.font.bold = True
            ans_p.font.color.rgb = DARK

            # Checkmark directly beside answer text, not diagonally above
            check_left = answer_box_left + answer_box_width + Inches(0.18)

            add_circle(
                slide2,
                check_left,
                row_top + Inches(0.24),
                check_size,
                GREEN,
                text="✓",
                text_rgb=WHITE,
                font_size=24,
            )

            exp_box = slide2.shapes.add_textbox(
                Inches(0.85),
                Inches(6.38),
                Inches(11.65),
                Inches(0.62),
            )
            exp_tf = exp_box.text_frame
            set_text_frame(exp_tf, 0, 0, 0, 0, wrap=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)

            exp_p = exp_tf.paragraphs[0]
            exp_p.text = explanation
            exp_p.alignment = PP_ALIGN.CENTER
            exp_p.font.name = "Arial"
            exp_p.font.size = fit_explanation_size(explanation)
            exp_p.font.italic = True
            exp_p.font.color.rgb = MID_GRAY

        print("[MCQ3] Saving...", flush=True)
        prs.save(output_path)
        print(f"[MCQ3] Done — {len(mcqs) * 2} slides", flush=True)
        return True

    except Exception as e:
        import traceback
        print(f"[MCQ3] FAILED: {e}", flush=True)
        print(traceback.format_exc(), flush=True)
        return False
    
def create_vba_template_presentation_mobile(mcqs, output_path):
    if not PPTX_AVAILABLE:
        return False

    try:
        prs = Presentation()
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(13.33)

        BLUE = PptRGBColor(100, 149, 237)
        BG = PptRGBColor(249, 248, 242)
        DARK = PptRGBColor(22, 22, 26)
        WHITE = PptRGBColor(255, 255, 255)

        LEFT_MARGIN = Inches(0.6)
        BOX_WIDTH = prs.slide_width - Inches(1.2)

        for i, mcq in enumerate(mcqs, start=1):
            options = ensure_four_options(mcq)

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = BG
            bg.line.fill.background()

            qbox_top = Inches(1.8)
            qbox_height = Inches(3.2)

            qbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                LEFT_MARGIN,
                qbox_top,
                BOX_WIDTH,
                qbox_height,
            )
            qbox.fill.solid()
            qbox.fill.fore_color.rgb = BG
            qbox.line.color.rgb = BLUE
            qbox.line.width = Pt(6)
            qbox.adjustments[0] = 0.08

            qtf = qbox.text_frame
            qtf.word_wrap = True
            qtf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            qtf.margin_left = Inches(0.35)
            qtf.margin_right = Inches(0.35)
            qtf.margin_top = Inches(0.35)
            qtf.margin_bottom = Inches(0.35)

            qp = qtf.paragraphs[0]
            qp.text = f"Q{i}: {mcq['question']}"
            qp.font.name = "Arial"
            qp.font.size = question_font_size(mcq["question"], large=32, medium=28, small=24, tiny=20)
            qp.font.bold = True
            qp.font.color.rgb = DARK
            qp.alignment = PP_ALIGN.CENTER

            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                LEFT_MARGIN - Inches(0.5),
                qbox_top - Inches(0.4),
                Inches(1.2),
                Inches(1.2),
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = BLUE
            badge.line.fill.background()

            bt = badge.text_frame
            bt.clear()
            bp = bt.paragraphs[0]
            bp.text = str(i)
            bp.font.name = "Arial"
            bp.font.size = Pt(34)
            bp.font.bold = True
            bp.font.color.rgb = WHITE
            bp.alignment = PP_ALIGN.CENTER

            start_y = qbox_top + qbox_height + Inches(0.8)
            option_height = Inches(1.05)
            option_spacing = Inches(1.45)

            for opt in options:
                obox = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    LEFT_MARGIN,
                    start_y,
                    BOX_WIDTH,
                    option_height,
                )
                obox.fill.solid()
                obox.fill.fore_color.rgb = BG
                obox.line.color.rgb = BLUE
                obox.line.width = Pt(6)
                obox.adjustments[0] = 0.08

                ot = obox.text_frame
                ot.word_wrap = True
                ot.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                ot.margin_left = Inches(0.3)
                ot.margin_right = Inches(0.3)

                op = ot.paragraphs[0]
                op.text = opt
                op.font.name = "Arial"
                op.font.size = option_font_size(opt, large=26, medium=22, small=18)
                op.font.italic = True
                op.font.color.rgb = DARK
                op.alignment = PP_ALIGN.CENTER

                start_y += option_spacing

            slide2 = prs.slides.add_slide(prs.slide_layouts[6])

            bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg2.fill.solid()
            bg2.fill.fore_color.rgb = BG
            bg2.line.fill.background()

            qbox2_top = Inches(1.8)

            qbox2 = slide2.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                LEFT_MARGIN,
                qbox2_top,
                BOX_WIDTH,
                Inches(2.6),
            )
            qbox2.fill.solid()
            qbox2.fill.fore_color.rgb = BG
            qbox2.line.color.rgb = BLUE
            qbox2.line.width = Pt(6)
            qbox2.adjustments[0] = 0.08

            qtf2 = qbox2.text_frame
            qtf2.word_wrap = True
            qtf2.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            qtf2.margin_left = Inches(0.35)
            qtf2.margin_right = Inches(0.35)

            qtf2.paragraphs[0].text = mcq["question"]
            qtf2.paragraphs[0].font.size = question_font_size(mcq["question"], large=28, medium=24, small=20, tiny=18)
            qtf2.paragraphs[0].font.bold = True
            qtf2.paragraphs[0].font.color.rgb = DARK
            qtf2.paragraphs[0].alignment = PP_ALIGN.CENTER

            badge2 = slide2.shapes.add_shape(
                MSO_SHAPE.OVAL,
                LEFT_MARGIN - Inches(0.5),
                qbox2_top - Inches(0.4),
                Inches(1.2),
                Inches(1.2),
            )
            badge2.fill.solid()
            badge2.fill.fore_color.rgb = BLUE
            badge2.line.fill.background()

            bt2 = badge2.text_frame
            bt2.clear()
            bp2 = bt2.paragraphs[0]
            bp2.text = str(i)
            bp2.font.size = Pt(34)
            bp2.font.bold = True
            bp2.font.color.rgb = WHITE
            bp2.alignment = PP_ALIGN.CENTER

            abox = slide2.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                LEFT_MARGIN,
                qbox2_top + Inches(3.6),
                BOX_WIDTH,
                Inches(4.2),
            )
            abox.fill.solid()
            abox.fill.fore_color.rgb = BG
            abox.line.color.rgb = BLUE
            abox.line.width = Pt(6)
            abox.adjustments[0] = 0.08

            atf = abox.text_frame
            atf.word_wrap = True
            atf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            atf.margin_left = Inches(0.35)
            atf.margin_right = Inches(0.35)
            atf.margin_top = Inches(0.35)

            p1 = atf.paragraphs[0]
            p1.text = f"Correct Answer: {mcq['answer']}"
            p1.font.size = Pt(30)
            p1.font.bold = True
            p1.font.color.rgb = DARK
            p1.alignment = PP_ALIGN.CENTER

            p2 = atf.add_paragraph()
            p2.text = mcq["explanation"]
            p2.font.size = question_font_size(mcq["explanation"], large=24, medium=22, small=20, tiny=18)
            p2.font.color.rgb = DARK
            p2.alignment = PP_ALIGN.CENTER

        prs.save(output_path)
        return True

    except Exception as e:
        print(f"Error creating mobile MCQ PPT: {e}")
        return False


def create_brand_template_presentation_from_ppt(drug_blocks, output_path):
    if not PPTX_AVAILABLE:
        return False

    try:
        template_path = "drug output template.pptx"
        prs = Presentation(template_path)
        template_slide = prs.slides[0]

        sldIdLst = list(prs.slides._sldIdLst)
        for sld in sldIdLst:
            rId = sld.rId
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(sld)

        def clone_template_slide():
            slide_layout = template_slide.slide_layout
            new_slide = prs.slides.add_slide(slide_layout)

            for shp in list(new_slide.shapes):
                el = shp._element
                el.getparent().remove(el)

            for shp in template_slide.shapes:
                new_el = deepcopy(shp._element)
                new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

            return new_slide

        for idx, drug in enumerate(drug_blocks, start=1):
            slide = clone_template_slide()
            shapes = slide.shapes

            generic = drug.get("generic", "")
            brand = drug.get("brand", "")
            header = drug.get("header", "")

            if len(shapes) >= 1:
                shapes[0].text = generic
            if len(shapes) >= 2:
                shapes[1].text = brand
            if len(shapes) >= 3:
                shapes[2].text = header
            if len(shapes) >= 4:
                shapes[3].text = str(idx)

        prs.save(output_path)
        return True

    except Exception as e:
        print(f"Error creating brand template PPT from PPT: {e}")
        return False


def create_brand_template_presentation(drug_blocks, output_path):
    if not PPTX_AVAILABLE:
        return False

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        BG = PptRGBColor(249, 244, 230)
        BLUE = PptRGBColor(51, 102, 204)
        YELLOW = PptRGBColor(255, 204, 0)
        BLACK = PptRGBColor(0, 0, 0)
        WHITE = PptRGBColor(255, 255, 255)
        GREY = PptRGBColor(120, 120, 120)

        for idx, drug in enumerate(drug_blocks, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            generic = drug.get("generic", "")
            brand = drug.get("brand", "")
            extra = drug.get("extra", "")

            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = BG
            bg.line.fill.background()

            blue_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5),
                Inches(1.3),
                Inches(12.3),
                Inches(1.2),
            )
            blue_box.fill.solid()
            blue_box.fill.fore_color.rgb = BLUE
            blue_box.line.fill.background()

            gtf = blue_box.text_frame
            gtf.word_wrap = True
            gtf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = gtf.paragraphs[0]
            p.text = generic
            p.font.name = "Arial"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

            brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.0))
            btf = brand_box.text_frame
            btf.word_wrap = True
            btf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            pb = btf.paragraphs[0]
            pb.text = brand
            pb.font.name = "Arial"
            pb.font.size = Pt(40)
            pb.font.bold = True
            pb.font.color.rgb = BLACK
            pb.alignment = PP_ALIGN.CENTER

            yellow_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5),
                Inches(4.4),
                Inches(12.3),
                Inches(1.6),
            )
            yellow_box.fill.solid()
            yellow_box.fill.fore_color.rgb = YELLOW
            yellow_box.line.fill.background()

            etf = yellow_box.text_frame
            etf.word_wrap = True
            etf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            etf.margin_left = Pt(20)
            etf.margin_right = Pt(20)
            pe = etf.paragraphs[0]
            pe.text = extra
            pe.font.name = "Arial"
            pe.font.size = Pt(32)
            pe.font.color.rgb = BLACK
            pe.alignment = PP_ALIGN.CENTER

            num = slide.shapes.add_textbox(
                prs.slide_width - Inches(1.0),
                prs.slide_height - Inches(0.9),
                Inches(0.8),
                Inches(0.8),
            )
            ntf = num.text_frame
            pn = ntf.paragraphs[0]
            pn.text = str(idx)
            pn.font.name = "Arial"
            pn.font.size = Pt(28)
            pn.font.color.rgb = GREY
            pn.alignment = PP_ALIGN.RIGHT

        prs.save(output_path)
        return True

    except Exception as e:
        print(f"Error in create_brand_template_presentation: {e}")
        return False


def create_brand_template_presentation_mobile(drug_blocks, output_path):
    if not PPTX_AVAILABLE:
        return False

    try:
        prs = Presentation()
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(13.33)

        BG = PptRGBColor(249, 244, 230)
        BLUE = PptRGBColor(51, 102, 204)
        YELLOW = PptRGBColor(255, 204, 0)
        BLACK = PptRGBColor(0, 0, 0)
        WHITE = PptRGBColor(255, 255, 255)

        for idx, drug in enumerate(drug_blocks, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            generic = drug.get("generic", "")
            brand = drug.get("brand", "")
            extra = drug.get("extra", "")

            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = BG
            bg.line.fill.background()

            generic_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5),
                Inches(1.0),
                prs.slide_width - Inches(1.0),
                Inches(1.8),
            )
            generic_box.fill.solid()
            generic_box.fill.fore_color.rgb = BLUE
            generic_box.line.fill.background()

            tf = generic_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]
            p.text = generic
            p.font.name = "Arial"
            p.font.size = Pt(42)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

            brand_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(3.4),
                prs.slide_width - Inches(1.0),
                Inches(2.0),
            )
            btf = brand_box.text_frame
            btf.word_wrap = True
            btf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            bp = btf.paragraphs[0]
            bp.text = brand
            bp.font.name = "Arial"
            bp.font.size = Pt(40)
            bp.font.bold = True
            bp.font.color.rgb = BLACK
            bp.alignment = PP_ALIGN.CENTER

            extra_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5),
                prs.slide_height - Inches(3.0),
                prs.slide_width - Inches(1.0),
                Inches(1.8),
            )
            extra_box.fill.solid()
            extra_box.fill.fore_color.rgb = YELLOW
            extra_box.line.fill.background()

            etf = extra_box.text_frame
            etf.word_wrap = True
            etf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            ep = etf.paragraphs[0]
            ep.text = extra
            ep.font.name = "Arial"
            ep.font.size = Pt(30)
            ep.font.color.rgb = BLACK
            ep.alignment = PP_ALIGN.CENTER

        prs.save(output_path)
        return True

    except Exception as e:
        print(f"Error creating mobile brand template: {e}")
        return False


def create_ppt_template_presentation(mcqs, output_path):
    if not PPTX_AVAILABLE:
        return False

    try:
        prs = Presentation("templates/ppt_template.pptx")

        layout = None
        for slide_layout in prs.slide_layouts:
            if slide_layout.placeholders:
                layout = slide_layout
                break

        if not layout:
            layout = prs.slide_layouts[0]

        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)

        for i, mcq in enumerate(mcqs, 1):
            slide = prs.slides.add_slide(layout)

            if slide.shapes.title:
                slide.shapes.title.text = f"Q{i}: {mcq['question']}"

            content_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:
                    content_shape = shape
                    break

            if content_shape:
                tf = content_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                for option in ensure_four_options(mcq):
                    p = tf.add_paragraph()
                    p.text = option
                    p.font.size = Pt(18)

                if mcq["answer"] or mcq["explanation"]:
                    p = tf.add_paragraph()
                    answer_text = f"Answer: {mcq['answer']}" if mcq["answer"] else ""
                    explanation_text = (
                        f"Explanation: {mcq['explanation']}" if mcq["explanation"] else ""
                    )
                    p.text = f"{answer_text}\n{explanation_text}".strip()
                    p.font.bold = True
                    p.font.size = Pt(16)

        prs.save(output_path)
        return True

    except Exception as e:
        print(f"Error creating PPT template presentation: {e}")
        return False

def translate_texts_to_french_batch(texts: list) -> list:
    import time as _time
 
    clean_texts = [t if t else "" for t in texts]
    if not any(t.strip() for t in clean_texts):
        return clean_texts
 
    joined = "\n---BLOCK-END---\n".join(clean_texts)
 
    for attempt in range(2):                              # ← CHANGED: 2 tries max
        try:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "Translate ALL English text into natural, professional French. "
                            "Do not skip short words, labels, or headings. "
                            "Preserve numbers, punctuation, and formatting markers. "
                            "Keep every block separated exactly with ---BLOCK-END---."
                        ),
                    },
                    {"role": "user", "content": joined},
                ],
                temperature=0,
                timeout=45,                              # ← CHANGED: hard timeout
            )
 
            translated = response.choices[0].message.content or joined
            parts = [p.strip() for p in translated.split("---BLOCK-END---")]
 
            if len(parts) == len(clean_texts):
                return parts
 
            # count mismatch — fall through to per-text fallback below
            break
 
        except Exception as e:
            print(f"Batch translate attempt {attempt+1} failed: {e}")
            if attempt == 0:
                _time.sleep(1)
 
    # Per-text fallback (only reached on mismatch or error)
    results = []
    for text in clean_texts:
        if not text.strip():
            results.append(text)
            continue
        try:
            r = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Translate to French. Return only the translation."},
                    {"role": "user", "content": text},
                ],
                temperature=0,
                timeout=15,                               # ← CHANGED
            )
            results.append(r.choices[0].message.content.strip())
        except Exception:
            results.append(text)   # keep original on failure
    return results
 

def translate_docx_embedded_images(doc):
    return
    image_ext_map = {
        "image/png": "png", "image/jpeg": "jpg", "image/jpg": "jpg",
        "image/bmp": "bmp", "image/tiff": "tiff", "image/gif": "png",
    }
 
    parts_to_check = [doc.part]
    for section in doc.sections:
        parts_to_check.append(section.header.part)
        parts_to_check.append(section.footer.part)
 
    image_jobs = []
    seen = set()
    for part in parts_to_check:
        for rel in part.rels.values():
            if "image" not in rel.reltype:
                continue

            if getattr(rel, "is_external", False):
                print("Skipping external image relationship")
                continue

            try:
                image_part = rel.target_part
            except Exception as e:
                print(f"Skipping unreadable image relationship: {e}")
                continue
            if id(image_part) in seen:
                continue
            seen.add(id(image_part))
            content_type = getattr(image_part, "content_type", "image/png")
            ext = image_ext_map.get(content_type, "png")
            image_jobs.append((image_part, ext))
 
    if not image_jobs:
        return
 
    def process_image(job):
        image_part, ext = job
        # ← CHANGED: uuid avoids any filename collision between threads
        uid = uuid.uuid4().hex
        temp_input  = os.path.join(app.config["IMAGES_FOLDER"], f"img_in_{uid}.{ext}")
        temp_output = os.path.join(app.config["IMAGES_FOLDER"], f"img_out_{uid}.{ext}")
        new_blob = None
 
        try:
            with open(temp_input, "wb") as f:
                f.write(image_part.blob)
 
            success, msg = translate_and_overlay_text(temp_input, temp_output, client=client)
            print("DOCX image:", msg)
 
            # ← CHANGED: read blob INSIDE the try, before finally cleans up
            if success and os.path.exists(temp_output):
                with open(temp_output, "rb") as f:
                    new_blob = f.read()
 
        except Exception as e:
            print(f"Embedded image translation skipped: {e}")
 
        finally:
            # ← CHANGED: delete files one at a time with individual guards
            for p in [temp_input, temp_output]:
                try:
                    if os.path.exists(p):
                        os.remove(p)
                except OSError as e:
                    print(f"Could not delete temp file {p}: {e}")
 
        # Return outside finally so the blob is never at risk
        if new_blob and len(new_blob) > 1000:
            return image_part, new_blob
        return None
 
    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = {executor.submit(process_image, job): job for job in image_jobs}
        for future in as_completed(futures):
            result = future.result()
            if result:
                image_part, new_blob = result
                image_part._blob = new_blob

def add_word_textbox_overlays_with_word(docx_path: str):
    print("Skipping Word COM overlay; using PIL embedded-image replacement instead.")
    return

def translate_docx_keep_layout(input_path: str, output_path: str):
    import time as _time
 
    doc = Document(input_path)
    paragraphs = []
 
    def collect_paragraphs_from_container(container):
        for paragraph in container.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(paragraph)
        for table in container.tables:
            for row in table.rows:
                for cell in row.cells:
                    collect_paragraphs_from_container(cell)
 
    collect_paragraphs_from_container(doc)
    for section in doc.sections:
        collect_paragraphs_from_container(section.header)
        collect_paragraphs_from_container(section.footer)

    # Collect Word text boxes / shapes
    from docx.text.paragraph import Paragraph

    for txbx in doc.element.xpath(".//*[local-name()='txbxContent']"):
        for p_el in txbx.xpath(".//*[local-name()='p']"):
            p = Paragraph(p_el, doc)
            if p.text.strip():
                paragraphs.append(p)

    batch_size = 80

    all_runs = []
    for paragraph in paragraphs:
        for run in paragraph.runs:
            if run.text and run.text.strip():
                if not run._element.xpath(".//*[local-name()='drawing']"):
                    all_runs.append(run)

    for i in range(0, len(all_runs), batch_size):
        batch = all_runs[i:i + batch_size]
        original_texts = [r.text for r in batch]

        start = time.time()
        translated_texts = translate_texts_to_french_batch(original_texts)
        print(f"Batch took {round(time.time() - start, 2)}s")

        for run, translated in zip(batch, translated_texts):
            run.text = translated
 
    print("NOW TRANSLATING EMBEDDED DOCX IMAGES (parallel)...")
    translate_docx_embedded_images(doc)                   # ← now parallel
 
    doc.save(output_path)
# ---------------- ROUTES ----------------
@app.route("/upload-image", methods=["POST"])
def upload_image():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400

    if file and allowed_image_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config["IMAGES_FOLDER"], filename)
        file.save(filepath)
        return jsonify({
            "message": "✅ Image uploaded successfully", 
            "filename": filename,
            "image_url": f"/images/{filename}"
        })

    return jsonify({"error": "Invalid image file type"}), 400

def allowed_image_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_IMAGE_EXTENSIONS

@app.route("/translate-image", methods=["POST"])
def translate_image():
    try:
        data = request.get_json(silent=True) or request.form
        filename = data.get("filename")

        if not filename:
            return jsonify({"error": "No filename provided"}), 400

        filepath = os.path.join(app.config["IMAGES_FOLDER"], filename)
        if not os.path.exists(filepath):
            return jsonify({"error": f"Image file not found: {filepath}"}), 404

        base_name = filename.rsplit(".", 1)[0]
        ext = filename.rsplit(".", 1)[1]
        output_filename = f"{base_name}_french_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"
        output_path = os.path.join(app.config["GENERATED_FOLDER"], output_filename)

        success, message = translate_and_overlay_text(filepath, output_path, client=client)

        if success:
            return jsonify({
                "message": f"✅ Image translated successfully: {message}",
                "download_url": f"/download/{output_filename}"
            })
        else:
            return jsonify({"error": f"Translation failed: {message}"}), 500

    except Exception as e:
        print(f"/translate-image error: {e}")
        return jsonify({"error": f"Generation failed: {type(e).__name__}: {str(e)}"}), 500


@app.route("/images/<filename>")
def serve_image(filename):
    return send_from_directory(app.config["IMAGES_FOLDER"], filename)


@app.route("/")
def home():
    return send_from_directory(".", "index.html")


@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(filepath)
        return jsonify({"message": "✅ File uploaded successfully", "filename": filename})

    return jsonify({"error": "Invalid file type"}), 400

def find_binary(*names):
    for name in names:
        found = shutil.which(name)
        if found:
            return found
        if os.path.exists(name):
            return name
    return None


def convert_docx_to_html_pixel_perfect(input_path: str, output_path: str) -> str:
    """
    True visual DOCX -> HTML converter.

    Pipeline:
    1. LibreOffice renders DOCX to PDF.
    2. Poppler pdftoppm renders each PDF page to PNG.
    3. HTML embeds those page PNGs as base64.

    This is not semantic/editable HTML. It is visual HTML.
    It is the closest option for 1:1 layout, colors, fonts, equations,
    bullets, tables, and rotated text.
    """

    import base64
    import subprocess
    import tempfile
    import shutil
    from pathlib import Path

    libreoffice = find_binary(
        "libreoffice",
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/libreoffice",
        "/usr/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    )

    pdftoppm = find_binary(
        "pdftoppm",
        r"C:\poppler-26.02.0\Library\bin\pdftoppm.exe",
        "/usr/bin/pdftoppm",
    )

    if not libreoffice:
        raise RuntimeError("LibreOffice is not installed or not found in PATH.")

    if not pdftoppm:
        raise RuntimeError("pdftoppm is not installed or not found in PATH.")

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)

        # 1. DOCX -> PDF using LibreOffice
        pdf_result = subprocess.run(
            [
                libreoffice,
                "--headless",
                "--norestore",
                "--nofirststartwizard",
                "--convert-to",
                "pdf",
                "--outdir",
                str(tmpdir_path),
                input_path,
            ],
            capture_output=True,
            text=True,
            timeout=180,
        )

        if pdf_result.returncode != 0:
            raise RuntimeError(
                "LibreOffice PDF conversion failed:\n"
                + (pdf_result.stderr or pdf_result.stdout or "Unknown error")
            )

        pdf_files = list(tmpdir_path.glob("*.pdf"))

        if not pdf_files:
            raise RuntimeError("LibreOffice did not create a PDF file.")

        pdf_path = pdf_files[0]

        # 2. PDF -> PNG pages
        # 220 DPI is a good quality/size balance.
        page_prefix = tmpdir_path / "page"

        image_result = subprocess.run(
        [
            pdftoppm,
            "-png",
            "-r",
            "300",
            str(pdf_path),
            str(page_prefix),
        ],
            capture_output=True,
            text=True,
            timeout=180,
        )

        if image_result.returncode != 0:
            raise RuntimeError(
                "pdftoppm page rendering failed:\n"
                + (image_result.stderr or image_result.stdout or "Unknown error")
            )

        page_images = sorted(tmpdir_path.glob("page-*.png"))

        if not page_images:
            raise RuntimeError("No PNG pages were created from the PDF.")

        # 3. Build self-contained HTML with base64 images
        page_html = []

        for idx, img_path in enumerate(page_images, start=1):
            with open(img_path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("utf-8")

            page_html.append(
                f"""
                <section class="page">
                  <img src="data:image/png;base64,{b64}" alt="Page {idx}" />
                </section>
                """
            )

        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8" />
<meta name="viewport" content="width=device-width, initial-scale=1.0" />
<title>Converted Document</title>
<style>
  * {{
    box-sizing: border-box;
  }}

  html,
  body {{
    margin: 0;
    padding: 0;
    background: #ffffff;
    font-family: Arial, sans-serif;
  }}

  .document {{
    width: 100%;
    max-width: 980px;
    margin: 0 auto;
    display: block;
    background: #ffffff;
  }}

  .page {{
    width: 100%;
    margin: 0;
    padding: 0;
    background: #ffffff;
    box-shadow: none;
    line-height: 0;
  }}

  .page img {{
    display: block;
    width: 100%;
    height: auto;
    margin: 0;
    padding: 0;
  }}

  @media print {{
    .page {{
      page-break-after: always;
    }}
  }}
</style>
</head>
<body>
  <main class="document">
    {''.join(page_html)}
  </main>
</body>
</html>
"""

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)

        return html
    
@app.route("/generated-assets/<asset_folder>/<filename>")
def generated_assets(asset_folder, filename):
    safe_folder = secure_filename(asset_folder)
    safe_file = secure_filename(filename)
    asset_dir = os.path.join(app.config["GENERATED_FOLDER"], safe_folder)
    return send_from_directory(asset_dir, safe_file)
@app.route("/generate", methods=["POST"])
def generate_output():
    import uuid as _uuid

    job_id = _uuid.uuid4().hex
    job_progress[job_id] = {"pct": 3, "state": "Received request", "eta": "Starting..."}

    def _set(pct, state, eta=""):
        job_progress[job_id] = {"pct": pct, "state": state, "eta": eta}

    try:
        data = request.get_json(silent=True) or request.form
        filename = data.get("filename")
        template_choice = data.get("template")
        cbt_topic = (data.get("cbt_topic") or "").strip()

        if not template_choice:
            return jsonify({"error": "Missing parameters"}), 400

        # ── CBT GENERATOR ──────────────────────────────────────────────────
        if template_choice == "cbt":
            if not cbt_topic:
                return jsonify({"error": "Please enter a topic for CBT Generator"}), 400
            _set(10, "Sending to GPT-4o", "~15 seconds")
            prompt = f"""
Generate 25 MCQs on the topic: "{cbt_topic}"
Format MUST match EXACTLY this structure:
**1.** *Question text here…*
a. Option A
b. Option B
c. Option C
d. Option D
**Answer: c**
**Explanation:** *Short explanation here.*
Rules:
- Same line breaks
- Numbered 1–25
- Keep formatting consistent
"""
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "system", "content": "You generate perfectly formatted medical MCQs."}, {"role": "user", "content": prompt}],
                temperature=0,
            )
            _set(80, "Building Word document", "~3 seconds")
            content = response.choices[0].message.content or ""
            output_filename = f"CBT_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            output_path = os.path.join(app.config["GENERATED_FOLDER"], output_filename)
            doc = Document()
            for line in content.split("\n"):
                line = line.rstrip()
                if line.startswith("**") and line.endswith("**"):
                    p = doc.add_paragraph(); r = p.add_run(line.replace("**", "")); r.bold = True; continue
                if line.startswith("**Answer"):
                    p = doc.add_paragraph(); r = p.add_run(line.replace("**", "")); r.bold = True; continue
                if line.startswith("**Explanation"):
                    p = doc.add_paragraph(); r = p.add_run(line.replace("**", "")); r.bold = True; r.italic = True; continue
                if line.startswith("*") and line.endswith("*"):
                    p = doc.add_paragraph(); r = p.add_run(line.replace("*", "")); r.italic = True; continue
                p = doc.add_paragraph(); p.add_run(line)
            doc.save(output_path)
            _set(100, "Done", "Complete")
            return jsonify({"message": "✅ CBT Word file generated!", "download_url": f"/download/{output_filename}", "job_id": job_id})

        # ── FILE-BASED MODES ───────────────────────────────────────────────
        if not filename:
            return jsonify({"error": "Missing filename"}), 400
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        if not os.path.exists(filepath):
            return jsonify({"error": f"Uploaded file not found: {filepath}"}), 404

        # ── CHAPTER TRANSLATE ──────────────────────────────────────────────
        if template_choice == "chapter_translate":
            ext = filename.rsplit(".", 1)[-1].lower()
            if ext != "docx":
                return jsonify({"error": "Chapter Translator only supports Word .docx files."}), 400
            _set(10, "Reading document", "~45 seconds")
            output_filename = f"translated_chapter_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            output_path = os.path.join(app.config["GENERATED_FOLDER"], output_filename)

            # Patch translate_texts_to_french_batch to report progress
            doc_obj = Document(filepath)
            paragraphs = []
            def _collect(container):
                for p in container.paragraphs:
                    if p.text.strip(): paragraphs.append(p)
                for t in container.tables:
                    for row in t.rows:
                        for cell in row.cells: _collect(cell)
            _collect(doc_obj)
            for section in doc_obj.sections:
                _collect(section.header); _collect(section.footer)
            from docx.text.paragraph import Paragraph as _Para
            for txbx in doc_obj.element.xpath(".//*[local-name()='txbxContent']"):
                for p_el in txbx.xpath(".//*[local-name()='p']"):
                    p = _Para(p_el, doc_obj)
                    if p.text.strip(): paragraphs.append(p)

            all_runs = []
            for paragraph in paragraphs:
                for run in paragraph.runs:
                    if run.text and run.text.strip():
                        if not run._element.xpath(".//*[local-name()='drawing']"):
                            all_runs.append(run)

            batch_size = 35
            total_batches = max(1, (len(all_runs) + batch_size - 1) // batch_size)
            for batch_num, i in enumerate(range(0, len(all_runs), batch_size)):
                pct = 10 + int((batch_num / total_batches) * 78)
                remaining_batches = total_batches - batch_num
                eta_secs = remaining_batches * 4
                _set(pct, f"Translating batch {batch_num + 1}/{total_batches}", f"~{eta_secs}s remaining")
                batch = all_runs[i:i + batch_size]
                translated_texts = translate_texts_to_french_batch([r.text for r in batch])
                for run, translated in zip(batch, translated_texts):
                    run.text = translated

            _set(92, "Saving document", "~2 seconds")
            translate_docx_embedded_images(doc_obj)
            doc_obj.save(output_path)
            _set(100, "Done", "Complete")
            return jsonify({"message": "✅ Word chapter translated to French", "download_url": f"/download/{output_filename}", "job_id": job_id})

        # ── DOCX TO HTML ───────────────────────────────────────────────────
        if template_choice == "docx_to_html":
            ext = filename.rsplit(".", 1)[-1].lower()
            if ext != "docx":
                return jsonify({"error": "HTML converter only supports .docx files"}), 400
            _set(20, "Converting with LibreOffice", "~8 seconds")
            output_filename = f"converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
            output_path = os.path.join(app.config["GENERATED_FOLDER"], output_filename)
            try:
                # Try the new high-fidelity converter first
                html_output = convert_docx_to_html_pixel_perfect(filepath, output_path)

            except Exception as e:
                print(f"Pixel-perfect converter failed: {type(e).__name__}: {e}")
                print("Falling back to old DOCX → HTML converter...")

                try:
                    # Old working converter fallback
                    html_output = convert_docx_to_html_libreoffice(filepath, output_path)

                except Exception as e2:
                    print(f"LibreOffice HTML fallback failed: {type(e2).__name__}: {e2}")
                    print("Falling back to manual parser...")

                    try:
                        html_output = convert_docx_to_html_faithful(filepath)
                        with open(output_path, "w", encoding="utf-8") as f:
                            f.write(html_output)

                    except Exception as e3:
                        return jsonify({
                            "error": (
                                "DOCX to HTML failed. "
                                f"Pixel-perfect: {type(e).__name__}: {str(e)} | "
                                f"LibreOffice fallback: {type(e2).__name__}: {str(e2)} | "
                                f"Manual fallback: {type(e3).__name__}: {str(e3)}"
                            )
                        }), 500
            _set(100, "Done", "Complete")
            return jsonify({
                "message": "✅ HTML generated successfully",
                "download_url": f"/download/{output_filename}",
                "copy_url": f"/raw-html/{output_filename}",
                "job_id": job_id
            })

        # ── All remaining modes need file text ─────────────────────────────
        _set(8, "Extracting text", "~2 seconds")
        file_content = extract_text(filepath)
        if not file_content.strip():
            return jsonify({"error": "File is empty"}), 400

        # ── BRAND / DRUG TEMPLATES ─────────────────────────────────────────
        if template_choice in {"brand_template", "brand_template_mobile"}:
            _set(20, "Parsing drug blocks", "~2 seconds")
            drug_blocks = parse_drug_brand_blocks(file_content)
            if not drug_blocks:
                return jsonify({"error": "No brand/generic blocks were found in this file."}), 400
            _set(55, "Building presentation", "~5 seconds")
            output_filename = f"generated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            output_path = os.path.join(app.config["GENERATED_FOLDER"], output_filename)
            success = create_brand_template_presentation(drug_blocks, output_path) if template_choice == "brand_template" else create_brand_template_presentation_mobile(drug_blocks, output_path)
            _set(100, "Done", "Complete")
            if success:
                return jsonify({"message": "File generated successfully", "download_url": f"/download/{output_filename}", "job_id": job_id})
            return jsonify({"error": "Failed to generate file"}), 500

        # ── MCQ EXTRACTION ─────────────────────────────────────────────────
        _set(15, "Normalizing MCQs with GPT", "~20 seconds")
        chunks = split_text_for_llm(file_content)
        all_mcqs_raw = []
        total_chunks = len(chunks)
        for chunk_num, chunk in enumerate(chunks):
            pct = 15 + int((chunk_num / total_chunks) * 55)
            _set(pct, f"Processing chunk {chunk_num + 1}/{total_chunks}", f"~{(total_chunks - chunk_num) * 5}s remaining")
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You normalize messy exam content into strict MCQ format."},
                    {"role": "user", "content": build_mcq_normalization_prompt(chunk)},
                ],
                temperature=0,
            )
            chunk_text = normalize_mcq_output_text(response.choices[0].message.content or "")
            chunk_mcqs = parse_normalized_mcq_text(chunk_text)
            all_mcqs_raw.extend(chunk_mcqs)

        normalized_mcq_text = render_normalized_mcq_text(all_mcqs_raw)
        mcqs = all_mcqs_raw

        if not mcqs:
            mcqs = regex_extract_mcqs_fallback(file_content)
            normalized_mcq_text = render_normalized_mcq_text(mcqs) if mcqs else ""

        if not mcqs:
            return jsonify({"error": "No MCQs could be extracted from this file."}), 400

        # ── CBT PARSER ─────────────────────────────────────────────────────
        if template_choice == "cbt_parser":
            _set(85, "Writing Word document", "~2 seconds")
            output_filename = f"CBT_PARSED_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            output_path = os.path.join(app.config["GENERATED_FOLDER"], output_filename)
            write_normalized_mcqs_to_docx(normalized_mcq_text, output_path)
            _set(100, "Done", "Complete")
            return jsonify({"message": "✅ File parsed and formatted successfully", "download_url": f"/download/{output_filename}", "job_id": job_id})

        # ── MCQ GRADER ─────────────────────────────────────────────────────
        if template_choice == "mcq_grader":
            grader_mcqs = [{"question": m["question"], "answer": m["answer"]} for m in mcqs if (m.get("answer") or "").upper() in {"A", "B", "C", "D"}]
            if not grader_mcqs:
                return jsonify({"error": "No MCQs with clear A–D answers were found."}), 400
            _set(100, "Done", "Complete")
            return jsonify({"message": "MCQs loaded for grading.", "mcqs": grader_mcqs, "job_id": job_id})

        # ── PPT GENERATION ─────────────────────────────────────────────────
        _set(75, "Building presentation", f"~{max(3, len(mcqs) // 5)}s remaining")
        output_filename = f"generated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        output_path = os.path.join(app.config["GENERATED_FOLDER"], output_filename)
        fn_map = {"vba": create_vba_template_presentation, "mcq2": create_mcq_generator2_exact, "mcq3": create_mcq_generator3_exact, "vba_mobile": create_vba_template_presentation_mobile, "ppt": create_ppt_template_presentation}
        if template_choice not in fn_map:
            return jsonify({"error": "Invalid template type"}), 400
        success = fn_map[template_choice](mcqs, output_path)
        _set(100, "Done", "Complete")
        if success:
            return jsonify({"message": "File generated successfully", "download_url": f"/download/{output_filename}", "job_id": job_id})
        return jsonify({"error": "Failed to generate file"}), 500

    except Exception as e:
        print(f"/generate error: {e}")
        job_progress[job_id] = {"pct": 0, "state": "Failed", "eta": str(e)}
        return jsonify({"error": f"Generation failed: {type(e).__name__}: {str(e)}"}), 500


@app.route("/raw-html/<filename>")
def raw_html_file(filename):
    safe_name = secure_filename(filename)
    if not safe_name.lower().endswith(".html"):
        return jsonify({"error": "Only HTML files can be served here"}), 400

    path = os.path.join(app.config["GENERATED_FOLDER"], safe_name)
    if not os.path.exists(path):
        return jsonify({"error": "HTML file not found"}), 404

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()

    return app.response_class(html, mimetype="text/plain; charset=utf-8")
@app.route("/download/<filename>")
def download_file(filename):
    return send_from_directory(app.config["GENERATED_FOLDER"], filename, as_attachment=True)

@app.route("/progress/<job_id>")
def progress_stream(job_id):
    def generate():
        last = -1
        for _ in range(600):  # max 60s polling
            info = job_progress.get(job_id)
            if info is None:
                yield f"data: {{}}\n\n"
                break
            if info["pct"] != last:
                import json
                last = info["pct"]
                yield f"data: {json.dumps(info)}\n\n"
            if info["pct"] >= 100:
                break
            import time as _t
            _t.sleep(0.1)
    return app.response_class(generate(), mimetype="text/event-stream",
                              headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
# ---------------- RUN APP ----------------
if __name__ == "__main__":
    if not PPTX_AVAILABLE:
        print("WARNING: python-pptx is not installed. Please run: pip install python-pptx")
    app.run(host="0.0.0.0", port=8502, debug=False, threaded=True)